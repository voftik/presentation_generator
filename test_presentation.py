#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Тест качества презентации - проверка изображений, геометрии, шрифтов и границ
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError as e:
    print(f"ОШИБКА: Не установлены зависимости: {e}")
    print("Запустите: pip install python-pptx")
    sys.exit(1)


class PresentationTester:
    """Тестер качества презентации"""
    
    def __init__(self, presentation_path: str):
        self.presentation_path = Path(presentation_path)
        if not self.presentation_path.exists():
            raise FileNotFoundError(f"Презентация не найдена: {self.presentation_path}")
        
        self.prs = Presentation(str(self.presentation_path))
        self.issues = []
        self.warnings = []
        
    def test_all(self):
        """Запускает все тесты"""
        print("🔍 ТЕСТИРОВАНИЕ КАЧЕСТВА ПРЕЗЕНТАЦИИ")
        print("=" * 60)
        print(f"📂 Файл: {self.presentation_path}")
        print(f"📊 Слайдов: {len(self.prs.slides)}")
        
        # Общая информация о презентации
        print(f"📐 Размер слайда: {self.prs.slide_width} x {self.prs.slide_height}")
        print()
        
        # Запускаем тесты
        self.test_images()
        self.test_geometry()
        self.test_fonts()
        self.test_text_boundaries()
        self.test_slide_layouts()
        
        # Выводим результаты
        self.print_results()
    
    def test_images(self):
        """Тест наличия и сохранения изображений"""
        print("🖼️  ТЕСТ ИЗОБРАЖЕНИЙ")
        print("-" * 30)
        
        total_images = 0
        slides_with_images = 0
        image_positions = []
        
        for i, slide in enumerate(self.prs.slides):
            slide_images = 0
            
            for shape in slide.shapes:
                # Проверяем разные типы изображений
                if hasattr(shape, 'image'):
                    slide_images += 1
                    image_positions.append({
                        'slide': i + 1,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    })
                elif 'Picture' in str(type(shape)):
                    slide_images += 1
                    image_positions.append({
                        'slide': i + 1,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    })
            
            if slide_images > 0:
                slides_with_images += 1
            total_images += slide_images
        
        print(f"  ✅ Всего изображений: {total_images}")
        print(f"  ✅ Слайдов с изображениями: {slides_with_images}/{len(self.prs.slides)}")
        
        # Проверяем консистентность позиций изображений
        if len(image_positions) >= 4:  # Минимум 2 изображения на 2 слайдах
            # Группируем изображения по слайдам и проверяем консистентность
            images_per_slide = {}
            for pos in image_positions:
                slide_num = pos['slide']
                if slide_num not in images_per_slide:
                    images_per_slide[slide_num] = []
                images_per_slide[slide_num].append(pos)
            
            # Проверяем, что на каждом слайде одинаковое количество изображений
            image_counts = [len(images) for images in images_per_slide.values()]
            if len(set(image_counts)) == 1:
                images_count = image_counts[0]
                print(f"  ✅ На каждом слайде {images_count} изображений")
                
                # Проверяем геометрию для каждой позиции изображения
                consistent_positions = True
                slide_numbers = sorted(images_per_slide.keys())
                
                if len(slide_numbers) >= 2:
                    first_slide_images = images_per_slide[slide_numbers[0]]
                    
                    for slide_num in slide_numbers[1:]:
                        current_slide_images = images_per_slide[slide_num]
                        
                        for i in range(len(first_slide_images)):
                            ref_img = first_slide_images[i]
                            cur_img = current_slide_images[i]
                            
                            if (abs(cur_img['left'] - ref_img['left']) > 100 or 
                                abs(cur_img['top'] - ref_img['top']) > 100 or
                                abs(cur_img['width'] - ref_img['width']) > 100 or
                                abs(cur_img['height'] - ref_img['height']) > 100):
                                consistent_positions = False
                                break
                        
                        if not consistent_positions:
                            break
                
                if consistent_positions:
                    print(f"  ✅ Геометрия изображений консистентна")
                    first_img = first_slide_images[0]
                    print(f"     Позиция первого: ({first_img['left']//914:.0f}, {first_img['top']//914:.0f}) дюймов")
                    print(f"     Размер: {first_img['width']//914:.0f}x{first_img['height']//914:.0f} дюймов")
                else:
                    self.warnings.append("Изображения имеют незначительные различия в геометрии")
                    print(f"  ⚠️  Геометрия изображений имеет небольшие различия")
            else:
                self.issues.append("Разное количество изображений на слайдах")
                print(f"  ❌ Разное количество изображений на слайдах")
        
        if total_images == 0:
            self.issues.append("В презентации НЕТ изображений")
        
        print()
    
    def test_geometry(self):
        """Тест сохранения геометрии элементов"""
        print("📐 ТЕСТ ГЕОМЕТРИИ")
        print("-" * 30)
        
        title_positions = []
        content_positions = []
        
        for i, slide in enumerate(self.prs.slides):
            # Проверяем позицию заголовка
            if slide.shapes.title:
                title_positions.append({
                    'slide': i + 1,
                    'left': slide.shapes.title.left,
                    'top': slide.shapes.title.top,
                    'width': slide.shapes.title.width,
                    'height': slide.shapes.title.height
                })
            
            # Проверяем позицию контента
            for shape in slide.shapes:
                if (hasattr(shape, 'text') and 
                    shape != slide.shapes.title and 
                    shape.text.strip()):
                    content_positions.append({
                        'slide': i + 1,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    })
                    break  # Берем только первый контентный блок
        
        # Проверяем консистентность заголовков по типам макетов
        section_titles = []
        content_titles = []
        
        for i, slide in enumerate(self.prs.slides):
            layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
            
            if slide.shapes.title:
                title_info = {
                    'slide': i + 1,
                    'left': slide.shapes.title.left,
                    'top': slide.shapes.title.top,
                    'layout': layout_name
                }
                
                if 'Section' in layout_name:
                    section_titles.append(title_info)
                else:
                    content_titles.append(title_info)
        
        # Проверяем Section Header заголовки
        if len(section_titles) > 1:
            first_section = section_titles[0]
            section_consistent = all(
                abs(pos['left'] - first_section['left']) <= 100 and
                abs(pos['top'] - first_section['top']) <= 100
                for pos in section_titles[1:]
            )
            
            if not section_consistent:
                self.issues.append("Section Header заголовки имеют разные позиции")
        
        # Проверяем Title and Content заголовки
        if len(content_titles) > 1:
            first_content = content_titles[0]
            content_consistent = all(
                abs(pos['left'] - first_content['left']) <= 100 and
                abs(pos['top'] - first_content['top']) <= 100
                for pos in content_titles[1:]
            )
            
            if content_consistent:
                print(f"  ✅ Позиции заголовков консистентны (по типам макетов)")
            else:
                self.issues.append("Title and Content заголовки имеют разные позиции")
                print(f"  ❌ Позиции заголовков НЕ консистентны")
        else:
            print(f"  ✅ Позиции заголовков консистентны")
        
        # Проверяем консистентность контента
        if len(content_positions) > 1:
            first_content = content_positions[0]
            content_consistent = True
            
            for pos in content_positions[1:]:
                if (abs(pos['left'] - first_content['left']) > 100 or 
                    abs(pos['top'] - first_content['top']) > 100):
                    content_consistent = False
                    break
            
            if content_consistent:
                print(f"  ✅ Позиции контента консистентны")
            else:
                self.issues.append("Контентные блоки имеют разные позиции на слайдах")
                print(f"  ❌ Позиции контента НЕ консистентны")
        
        print()
    
    def test_fonts(self):
        """Тест правильности шрифтов"""
        print("🔤 ТЕСТ ШРИФТОВ")
        print("-" * 30)
        
        font_issues = []
        expected_font = 'Montserrat'
        
        for i, slide in enumerate(self.prs.slides):
            slide_layout = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
            
            # Проверяем заголовки
            if slide.shapes.title and hasattr(slide.shapes.title, 'text_frame'):
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_name = run.font.name
                        font_size = run.font.size
                        
                        if font_name and font_name != expected_font:
                            font_issues.append(f"Слайд {i+1}: заголовок использует {font_name} вместо {expected_font}")
                        
                        if font_size:
                            expected_title_size = Pt(36) if 'Section' in slide_layout else Pt(30)
                            if abs(font_size - expected_title_size) > Pt(2):
                                size_pt = int(font_size.pt) if hasattr(font_size, 'pt') else font_size
                                expected_pt = int(expected_title_size.pt)
                                font_issues.append(f"Слайд {i+1}: размер заголовка {size_pt}pt вместо {expected_pt}pt")
            
            # Проверяем контент
            for shape in slide.shapes:
                if (hasattr(shape, 'text_frame') and 
                    shape != slide.shapes.title and 
                    hasattr(shape, 'text') and shape.text.strip()):
                    
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font_name = run.font.name
                            font_size = run.font.size
                            
                            if font_name and font_name != expected_font:
                                font_issues.append(f"Слайд {i+1}: контент использует {font_name} вместо {expected_font}")
                            
                            if font_size:
                                expected_content_size = Pt(18)
                                if abs(font_size - expected_content_size) > Pt(2):
                                    size_pt = int(font_size.pt) if hasattr(font_size, 'pt') else font_size
                                    font_issues.append(f"Слайд {i+1}: размер контента {size_pt}pt вместо 18pt")
                    break  # Проверяем только первый контентный блок
        
        if font_issues:
            print(f"  ❌ Найдено {len(font_issues)} проблем со шрифтами:")
            for issue in font_issues[:5]:  # Показываем первые 5
                print(f"     • {issue}")
            if len(font_issues) > 5:
                print(f"     ... и еще {len(font_issues) - 5} проблем")
            self.issues.extend(font_issues)
        else:
            print(f"  ✅ Все шрифты корректны ({expected_font})")
            print(f"  ✅ Размеры шрифтов правильные")
        
        print()
    
    def test_text_boundaries(self):
        """Тест выхода текста за границы слайда"""
        print("📏 ТЕСТ ГРАНИЦ ТЕКСТА")
        print("-" * 30)
        
        boundary_issues = []
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        for i, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    # Проверяем границы shape
                    right_edge = shape.left + shape.width
                    bottom_edge = shape.top + shape.height
                    
                    if shape.left < 0:
                        boundary_issues.append(f"Слайд {i+1}: текст выходит за левую границу")
                    
                    if shape.top < 0:
                        boundary_issues.append(f"Слайд {i+1}: текст выходит за верхнюю границу")
                    
                    if right_edge > slide_width:
                        boundary_issues.append(f"Слайд {i+1}: текст выходит за правую границу")
                    
                    if bottom_edge > slide_height:
                        boundary_issues.append(f"Слайд {i+1}: текст выходит за нижнюю границу")
                    
                    # Проверяем длину текста
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            line_length = len(paragraph.text)
                            if line_length > 200:  # Очень длинная строка
                                self.warnings.append(f"Слайд {i+1}: очень длинная строка ({line_length} символов)")
        
        if boundary_issues:
            print(f"  ❌ Найдено {len(boundary_issues)} проблем с границами:")
            for issue in boundary_issues:
                print(f"     • {issue}")
            self.issues.extend(boundary_issues)
        else:
            print(f"  ✅ Весь текст находится в пределах слайда")
        
        print()
    
    def test_slide_layouts(self):
        """Тест правильности макетов слайдов"""
        print("🎨 ТЕСТ МАКЕТОВ")
        print("-" * 30)
        
        layout_stats = {}
        layout_issues = []
        
        for i, slide in enumerate(self.prs.slides):
            layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
            
            if layout_name not in layout_stats:
                layout_stats[layout_name] = 0
            layout_stats[layout_name] += 1
            
            # Проверяем соответствие макета содержимому
            title_text = slide.shapes.title.text if slide.shapes.title else ""
            
            # Очень простая эвристика для проверки
            if 'Section' in layout_name and len(title_text) > 80:
                layout_issues.append(f"Слайд {i+1}: слишком длинный заголовок для Section Header")
        
        print(f"  📊 Использованные макеты:")
        for layout, count in layout_stats.items():
            print(f"     • {layout}: {count} слайдов")
        
        expected_layouts = ['Section Header', 'Title and Content']
        unexpected_layouts = [l for l in layout_stats.keys() if l not in expected_layouts]
        
        if unexpected_layouts:
            self.warnings.extend([f"Неожиданный макет: {layout}" for layout in unexpected_layouts])
        
        if layout_issues:
            print(f"  ⚠️  Проблемы с макетами:")
            for issue in layout_issues:
                print(f"     • {issue}")
            self.warnings.extend(layout_issues)
        else:
            print(f"  ✅ Макеты используются корректно")
        
        print()
    
    def print_results(self):
        """Выводит итоговые результаты тестирования"""
        print("📋 ИТОГИ ТЕСТИРОВАНИЯ")
        print("=" * 60)
        
        # Критические проблемы
        if self.issues:
            print(f"❌ КРИТИЧЕСКИЕ ПРОБЛЕМЫ ({len(self.issues)}):")
            for issue in self.issues:
                print(f"   • {issue}")
            print()
        
        # Предупреждения
        if self.warnings:
            print(f"⚠️  ПРЕДУПРЕЖДЕНИЯ ({len(self.warnings)}):")
            for warning in self.warnings:
                print(f"   • {warning}")
            print()
        
        # Общая оценка
        if not self.issues and not self.warnings:
            print("🎉 ОТЛИЧНО! Презентация прошла все тесты без проблем")
            grade = "A+"
        elif not self.issues and len(self.warnings) <= 2:
            print("✅ ХОРОШО! Презентация качественная, есть незначительные замечания")
            grade = "A"
        elif not self.issues and len(self.warnings) <= 5:
            print("👍 УДОВЛЕТВОРИТЕЛЬНО! Презентация приемлемая, но нужны улучшения")
            grade = "B"
        elif len(self.issues) <= 3:
            print("⚠️  ТРЕБУЕТ ДОРАБОТКИ! Есть проблемы, которые нужно исправить")
            grade = "C"
        else:
            print("❌ НЕУДОВЛЕТВОРИТЕЛЬНО! Множественные критические проблемы")
            grade = "F"
        
        print(f"🏆 ИТОГОВАЯ ОЦЕНКА: {grade}")
        print("=" * 60)


def main():
    """Главная функция"""
    script_dir = Path(__file__).parent
    presentation_path = script_dir / "result" / "presentation.pptx"
    
    print("🔍 ТЕСТЕР КАЧЕСТВА ПРЕЗЕНТАЦИИ")
    print(f"📂 Проверяемый файл: {presentation_path}")
    print()
    
    try:
        tester = PresentationTester(str(presentation_path))
        tester.test_all()
    except FileNotFoundError as e:
        print(f"❌ ОШИБКА: {e}")
        print("💡 Сначала запустите: python main.py")
        sys.exit(1)
    except Exception as e:
        print(f"❌ ОШИБКА ТЕСТИРОВАНИЯ: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()