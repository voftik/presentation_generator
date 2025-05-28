#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ADHOC PowerPoint Generator v3 - с полным сохранением дизайна
Использует Markdown файл и шаблон с сохранением ВСЕХ элементов
H1 = Section Header слайд, H2 = Title and Content слайд
"""

import os
import sys
import re
from typing import List, Tuple, Dict, Optional
from pathlib import Path
import shutil

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.parts.image import ImagePart
except ImportError as e:
    print(f"ОШИБКА: Не установлены зависимости: {e}")
    print("Запустите: pip install python-pptx")
    sys.exit(1)


class MarkdownSlide:
    """Структура данных для слайда"""
    def __init__(self, slide_type: str, title: str, content: str = ""):
        self.slide_type = slide_type  # "section" или "content"
        self.title = title
        self.content = content


class AdvancedPowerPointGenerator:
    """Генератор презентаций с полным сохранением дизайна"""
    
    def __init__(self, markdown_path: str, template_path: str, output_dir: str):
        self.markdown_path = Path(markdown_path)
        self.template_path = Path(template_path)
        self.output_dir = Path(output_dir)
        
        # Создаем выходную директорию
        self.output_dir.mkdir(exist_ok=True)
        
        # Проверяем входные файлы
        if not self.markdown_path.exists():
            raise FileNotFoundError(f"Markdown файл не найден: {self.markdown_path}")
        if not self.template_path.exists():
            raise FileNotFoundError(f"Шаблон не найден: {self.template_path}")
    
    def parse_markdown(self) -> List[MarkdownSlide]:
        """Парсит Markdown файл в структуру слайдов"""
        print(f"📄 Читаем Markdown файл: {self.markdown_path}")
        
        with open(self.markdown_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        lines = content.split('\n')
        slides = []
        current_slide = None
        current_content_lines = []
        
        for line in lines:
            line = line.rstrip()
            
            # H1 заголовок - новый section слайд
            if line.startswith('# '):
                # Сохраняем предыдущий слайд
                if current_slide:
                    current_slide.content = '\n'.join(current_content_lines).strip()
                    slides.append(current_slide)
                
                # Создаем новый section слайд
                title = line[2:].strip()
                current_slide = MarkdownSlide("section", title)
                current_content_lines = []
            
            # H2 заголовок - новый content слайд
            elif line.startswith('## '):
                # Сохраняем предыдущий слайд
                if current_slide:
                    current_slide.content = '\n'.join(current_content_lines).strip()
                    slides.append(current_slide)
                
                # Создаем новый content слайд
                title = line[3:].strip()
                current_slide = MarkdownSlide("content", title)
                current_content_lines = []
            
            # H3 заголовки и обычный текст - добавляем к контенту
            elif line.startswith('### '):
                # H3 становится подзаголовком в контенте
                subtitle = line[4:].strip()
                current_content_lines.append(f"**{subtitle}**")
                current_content_lines.append("")  # Пустая строка после подзаголовка
            
            elif line.strip():
                # Обычный текст
                current_content_lines.append(line)
            
            elif current_content_lines and current_content_lines[-1]:  
                # Пустая строка (только если предыдущая не пустая)
                current_content_lines.append("")
        
        # Не забываем последний слайд
        if current_slide:
            current_slide.content = '\n'.join(current_content_lines).strip()
            slides.append(current_slide)
        
        print(f"✅ Обработано {len(slides)} слайдов")
        
        # Выводим структуру для проверки
        for i, slide in enumerate(slides[:5]):  # Первые 5 слайдов
            print(f"  Слайд {i+1} ({slide.slide_type}): {slide.title[:50]}...")
        
        return slides
    
    def create_presentation_from_template(self, slides: List[MarkdownSlide]) -> None:
        """Создает презентацию используя оригинальный шаблон как основу"""
        print("🎨 Создаем презентацию с сохранением дизайна...")
        
        try:
            # Загружаем шаблон напрямую без копирования
            prs = Presentation(str(self.template_path))
            
            # Сохраняем оригинальный слайд как эталон для стилей
            original_slide = prs.slides[0]
            
            # Получаем все изображения из оригинального слайда
            original_images = []
            for shape in original_slide.shapes:
                if hasattr(shape, 'image'):
                    original_images.append({
                        'shape': shape,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'image_data': shape.image.blob
                    })
            
            print(f"📷 Найдено {len(original_images)} изображений в шаблоне")
            
            # Очищаем все слайды
            slide_ids_to_remove = list(range(len(prs.slides)))
            for slide_id in reversed(slide_ids_to_remove):
                rId = prs.slides._sldIdLst[slide_id].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[slide_id]
            
            # Создаем слайды на основе разобранного Markdown
            slide_count = 0
            
            for slide_data in slides:
                if slide_data.slide_type == "section":
                    # Используем Section Header макет для H1
                    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section Header
                    self._update_section_slide(slide, slide_data.title)
                else:
                    # Используем Title and Content макет для H2
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
                    self._update_content_slide(slide, slide_data.title, slide_data.content)
                
                # Добавляем изображения на каждый слайд
                self._add_template_images(slide, original_images, prs)
                
                slide_count += 1
            
            print(f"✅ Создано {slide_count} слайдов")
            
            # Сохраняем презентацию
            output_path = self.output_dir / "presentation.pptx"
            prs.save(str(output_path))
            print(f"🎉 Презентация готова: {output_path}")
            
        except Exception as e:
            print(f"❌ Ошибка при создании презентации: {e}")
            import traceback
            traceback.print_exc()
            raise
    
    def _add_template_images(self, slide, original_images, presentation):
        """Добавляет изображения из шаблона на слайд"""
        try:
            from io import BytesIO
            
            for img_info in original_images:
                # Создаем изображение из blob данных
                image_stream = BytesIO(img_info['image_data'])
                
                # Добавляем изображение на слайд
                picture = slide.shapes.add_picture(
                    image_stream,
                    img_info['left'],
                    img_info['top'],
                    img_info['width'],
                    img_info['height']
                )
                
        except Exception as e:
            print(f"⚠️  Не удалось добавить изображения: {e}")
            # Продолжаем без изображений
    
    def _update_section_slide(self, slide, title: str):
        """Обновляет section слайд (H1)"""
        try:
            # Для Section Header макета
            if slide.shapes.title:
                slide.shapes.title.text = title
                self._format_section_title(slide.shapes.title)
                
        except Exception as e:
            print(f"⚠️  Проблема с section слайдом: {e}")
            # Fallback
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    shape.text = title
                    self._format_section_title(shape)
                    break
    
    def _update_content_slide(self, slide, title: str, content: str):
        """Обновляет content слайд (H2)"""
        try:
            # Заголовок
            if slide.shapes.title:
                slide.shapes.title.text = title
                self._format_content_title(slide.shapes.title)
            
            # Контент
            content_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_shape = shape
                    break
            
            if content_shape and content:
                formatted_content = self._format_markdown_content(content)
                content_shape.text = formatted_content
                self._format_content_text(content_shape)
                
        except Exception as e:
            print(f"⚠️  Проблема с content слайдом: {e}")
            # Fallback
            text_shapes = [s for s in slide.shapes if hasattr(s, 'text')]
            if len(text_shapes) >= 1:
                text_shapes[0].text = title
                self._format_content_title(text_shapes[0])
            if len(text_shapes) >= 2 and content:
                formatted_content = self._format_markdown_content(content)
                text_shapes[1].text = formatted_content
                self._format_content_text(text_shapes[1])
    
    def _format_markdown_content(self, content: str) -> str:
        """Форматирует Markdown контент для PowerPoint"""
        if not content:
            return ""
        
        lines = content.split('\n')
        formatted_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Убираем markdown разметку и форматируем
            if line.startswith('• '):
                # Уже bullet point
                formatted_lines.append(line)
            elif line.startswith('**') and line.endswith('**'):
                # Жирный текст - убираем звездочки, это будет подзаголовок
                clean_line = line[2:-2]
                formatted_lines.append(clean_line)
            elif line.startswith('*') or line.startswith('-'):
                # Markdown bullet points
                clean_line = line[1:].strip()
                formatted_lines.append(f"• {clean_line}")
            else:
                # Обычный текст
                if len(line) > 80:
                    # Длинный текст - оставляем как есть
                    formatted_lines.append(line)
                else:
                    # Короткий текст - делаем bullet point
                    formatted_lines.append(f"• {line}")
        
        return '\\n'.join(formatted_lines)
    
    def _format_section_title(self, shape):
        """Форматирует заголовок section слайда"""
        if not hasattr(shape, 'text_frame'):
            return
            
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'Montserrat'
                run.font.size = Pt(36)
                run.font.bold = True
    
    def _format_content_title(self, shape):
        """Форматирует заголовок content слайда"""
        if not hasattr(shape, 'text_frame'):
            return
            
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = 'Montserrat'
                run.font.size = Pt(30)
                run.font.bold = True
    
    def _format_content_text(self, shape):
        """Форматирует текст content слайда"""
        if not hasattr(shape, 'text_frame'):
            return
            
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = 'Montserrat'
                run.font.size = Pt(18)
                run.font.bold = False
    
    def generate(self):
        """Основной метод генерации"""
        print("🚀 Запуск генератора PowerPoint v3 (полное сохранение дизайна)")
        print("=" * 70)
        
        try:
            # Парсим Markdown
            slides = self.parse_markdown()
            
            # Создаем презентацию
            self.create_presentation_from_template(slides)
            
            print("=" * 70)
            print("✅ ГЕНЕРАЦИЯ ЗАВЕРШЕНА УСПЕШНО!")
            
        except Exception as e:
            print(f"❌ ОШИБКА: {e}")
            import traceback
            traceback.print_exc()
            sys.exit(1)


def main():
    """Главная функция"""
    script_dir = Path(__file__).parent
    
    # Пути к файлам
    markdown_path = script_dir / "content" / "content.md"
    template_path = script_dir / "tempate" / "Шаблон презентации 16х9.pptx"
    output_dir = script_dir / "result"
    
    print("🎯 ADHOC PowerPoint Generator v3")
    print(f"📂 Markdown файл: {markdown_path}")
    print(f"📂 Шаблон: {template_path}")
    print(f"📂 Результат: {output_dir}")
    print()
    
    # Создаем генератор и запускаем
    generator = AdvancedPowerPointGenerator(
        str(markdown_path), 
        str(template_path), 
        str(output_dir)
    )
    generator.generate()


if __name__ == "__main__":
    main()