"""
AI Image Generator для PowerPoint презентаций
Использует Gemini + Imagen для создания иллюстраций к слайдам
"""

import os
import re
import time
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    from google import genai
    from google.genai import types
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    print(f"❌ Необходимо установить зависимости: {e}")
    print("pip install google-genai pillow python-pptx")
    exit(1)


class PresentationImageGenerator:
    """Генератор изображений для презентаций с помощью Gemini + Imagen"""
    
    def __init__(self, gemini_api_key: str):
        self.client = genai.Client(api_key=gemini_api_key)
        self.model = "imagen-3.0-generate-002"  # Рабочая модель Imagen 3
        
        # Настройки для деловых презентаций
        self.business_style_prompts = {
            "corporate": "professional business corporate style, clean minimal design, blue and white colors",
            "technology": "modern technology style, digital abstract patterns, sleek and professional",
            "finance": "financial business style, charts and graphs aesthetic, professional corporate colors",
            "security": "cybersecurity style, network patterns, professional dark blue and orange",
            "innovation": "innovation and growth style, forward-thinking design, modern professional"
        }
        
    def analyze_slide_content(self, title: str, content: str) -> Dict[str, str]:
        """Анализирует содержимое слайда и определяет подходящий стиль изображения"""
        
        # Ключевые слова для определения тематики
        keywords_mapping = {
            "technology": ["технология", "цифровой", "автоматизация", "система", "программа", "данные"],
            "security": ["безопасность", "защита", "контроль", "мониторинг", "риск", "угроза"],
            "finance": ["финансы", "бюджет", "стоимость", "экономия", "деньги", "расходы"],
            "innovation": ["инновация", "развитие", "новый", "улучшение", "модернизация"],
            "corporate": ["управление", "организация", "процесс", "структура", "команда"]
        }
        
        text = (title + " " + content).lower()
        
        # Определяем стиль по ключевым словам
        style_scores = {}
        for style, keywords in keywords_mapping.items():
            score = sum(1 for keyword in keywords if keyword in text)
            style_scores[style] = score
        
        # Выбираем стиль с максимальным скором
        best_style = max(style_scores, key=style_scores.get) if max(style_scores.values()) > 0 else "corporate"
        
        return {
            "style": best_style,
            "style_prompt": self.business_style_prompts[best_style],
            "confidence": style_scores[best_style]
        }
    
    def generate_prompt_for_slide(self, title: str, content: str, slide_number: int) -> str:
        """Создает промпт для генерации изображения на основе содержимого слайда"""
        
        analysis = self.analyze_slide_content(title, content)
        
        # Базовый промпт для деловой презентации
        base_prompt = f"""
        Abstract professional illustration for business presentation slide,
        {analysis['style_prompt']},
        suitable for corporate presentation background,
        16:9 aspect ratio, high quality, no text or logos,
        subtle and elegant design that doesn't distract from content
        """
        
        # Добавляем контекстные элементы на основе содержимого
        if "данные" in content.lower() or "статистика" in content.lower():
            base_prompt += ", data visualization elements, abstract charts and graphs"
        elif "процесс" in content.lower() or "этап" in content.lower():
            base_prompt += ", workflow and process visualization, connected elements"
        elif "команда" in content.lower() or "сотрудник" in content.lower():
            base_prompt += ", teamwork and collaboration visualization, connected people silhouettes"
        elif "результат" in content.lower() or "достижение" in content.lower():
            base_prompt += ", success and achievement visualization, upward trending elements"
        
        return base_prompt.strip()
    
    def generate_image(self, prompt: str, output_path: str) -> bool:
        """Генерирует изображение по промпту и сохраняет в файл"""
        
        try:
            print(f"🎨 Генерируем изображение...")
            print(f"📝 Промпт: {prompt[:100]}...")
            
            response = self.client.models.generate_images(
                model=self.model,
                prompt=prompt,
                config=types.GenerateImagesConfig(
                    number_of_images=1,
                    aspect_ratio="16:9",
                    # safety_settings=types.SafetySetting.BLOCK_ONLY_HIGH  # Более мягкие фильтры для деловых изображений
                ),
            )
            
            # Сохраняем изображение
            img_bytes = response.generated_images[0].image.image_bytes
            Path(output_path).write_bytes(img_bytes)
            
            print(f"✅ Изображение сохранено: {output_path}")
            return True
            
        except Exception as e:
            print(f"❌ Ошибка генерации изображения: {e}")
            return False
    
    def add_image_to_slide(self, presentation_path: str, slide_number: int, image_path: str) -> bool:
        """Добавляет изображение на указанный слайд презентации"""
        
        try:
            # Открываем презентацию
            prs = Presentation(presentation_path)
            
            if slide_number >= len(prs.slides):
                print(f"❌ Слайд {slide_number + 1} не найден в презентации")
                return False
            
            slide = prs.slides[slide_number]
            
            # Определяем размеры для изображения (правый нижний угол)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_width = Inches(4)  # 4 дюйма ширина
            img_height = Inches(2.25)  # 2.25 дюйма высота (16:9 пропорция)
            
            # Позиционируем в правом нижнем углу
            left = slide_width - img_width - Inches(0.5)
            top = slide_height - img_height - Inches(0.5)
            
            # Добавляем изображение
            slide.shapes.add_picture(image_path, left, top, img_width, img_height)
            
            # Сохраняем презентацию
            prs.save(presentation_path)
            
            print(f"✅ Изображение добавлено на слайд {slide_number + 1}")
            return True
            
        except Exception as e:
            print(f"❌ Ошибка добавления изображения: {e}")
            return False
    
    def enhance_presentation_with_images(self, presentation_path: str, target_slides: List[int] = None) -> Dict:
        """Добавляет изображения к указанным слайдам презентации"""
        
        results = {
            "total_slides": 0,
            "generated_images": 0,
            "added_images": 0,
            "errors": []
        }
        
        try:
            # Открываем презентацию для анализа
            prs = Presentation(presentation_path)
            results["total_slides"] = len(prs.slides)
            
            # Если не указаны конкретные слайды, берем первые 5 (исключая титульный)
            if target_slides is None:
                target_slides = list(range(1, min(6, len(prs.slides))))  # Слайды 2-6
            
            print(f"🖼️ Генерируем изображения для слайдов: {[s+1 for s in target_slides]}")
            
            for slide_idx in target_slides:
                if slide_idx >= len(prs.slides):
                    continue
                
                slide = prs.slides[slide_idx]
                
                # Извлекаем текст слайда
                title = ""
                content = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        if shape.placeholder_format and shape.placeholder_format.type == 1:  # Title
                            title = text
                        else:
                            content += text + " "
                
                if not title and not content:
                    print(f"⚠️ Слайд {slide_idx + 1}: нет текста для анализа")
                    continue
                
                print(f"\n🔍 Анализируем слайд {slide_idx + 1}")
                print(f"📋 Заголовок: {title[:50]}...")
                print(f"📄 Содержимое: {content[:100]}...")
                
                # Генерируем промпт и изображение
                prompt = self.generate_prompt_for_slide(title, content, slide_idx)
                image_path = f"generated_image_slide_{slide_idx + 1}.png"
                
                if self.generate_image(prompt, image_path):
                    results["generated_images"] += 1
                    
                    # Добавляем изображение на слайд
                    if self.add_image_to_slide(presentation_path, slide_idx, image_path):
                        results["added_images"] += 1
                    else:
                        results["errors"].append(f"Не удалось добавить изображение на слайд {slide_idx + 1}")
                else:
                    results["errors"].append(f"Не удалось сгенерировать изображение для слайда {slide_idx + 1}")
                
                # Небольшая пауза между запросами
                time.sleep(2)
            
            print(f"\n📊 РЕЗУЛЬТАТЫ ГЕНЕРАЦИИ ИЗОБРАЖЕНИЙ:")
            print(f"Всего слайдов: {results['total_slides']}")
            print(f"Сгенерировано изображений: {results['generated_images']}")
            print(f"Добавлено в презентацию: {results['added_images']}")
            
            if results['errors']:
                print(f"⚠️ Ошибки ({len(results['errors'])}):")
                for error in results['errors']:
                    print(f"  - {error}")
            
            return results
            
        except Exception as e:
            print(f"❌ Критическая ошибка: {e}")
            results['errors'].append(str(e))
            return results


def main():
    """Тестирование генератора изображений"""
    
    # API ключ
    api_key = os.getenv('GEMINI_API_KEY') or "AIzaSyAx3AtOdrqJZ1fnqIE8fSrfiQWS3nHTs2I"
    if not api_key:
        print("❌ GEMINI_API_KEY environment variable is required")
        return
    
    # Создаем генератор
    generator = PresentationImageGenerator(api_key)
    
    # Тестируем на существующей презентации
    presentation_file = "result/presentation.pptx"
    
    if Path(presentation_file).exists():
        print("🎨 ГЕНЕРАТОР ИЗОБРАЖЕНИЙ ДЛЯ ПРЕЗЕНТАЦИЙ")
        print("=" * 60)
        
        # Генерируем изображения для слайдов 2, 3, 4
        results = generator.enhance_presentation_with_images(
            presentation_file, 
            target_slides=[1, 2, 3]  # Слайды 2, 3, 4 (индексы 1, 2, 3)
        )
        
        if results['added_images'] > 0:
            print(f"\n✅ УСПЕХ! Добавлено {results['added_images']} изображений")
            print(f"📁 Обновленная презентация: {presentation_file}")
        else:
            print("\n❌ Изображения не были добавлены")
            
    else:
        print(f"❌ Файл презентации не найден: {presentation_file}")
        print("Сначала создайте презентацию с помощью smart_generator.py")


if __name__ == "__main__":
    main()