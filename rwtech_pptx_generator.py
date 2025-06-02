#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RW Tech PPTX Generator - Universal PowerPoint Creator
Универсальный генератор презентаций PowerPoint от RW Tech
🚀 Advanced AI-Powered Presentation Generation System 🚀
"""

import os
import re
import sys
import json
import base64
import time
import logging
import threading
import queue
from datetime import datetime
from io import BytesIO
import requests
from tqdm import tqdm
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


class ProgressBar:
    """Анимированный прогресс-бар"""
    
    def __init__(self, total, prefix='', suffix='', decimals=1, length=50, fill='█', print_end="\r"):
        self.total = total
        self.prefix = prefix
        self.suffix = suffix
        self.decimals = decimals
        self.length = length
        self.fill = fill
        self.print_end = print_end
        self.current = 0
        
    def update(self, iteration):
        """Обновляет прогресс-бар"""
        self.current = iteration
        percent = ("{0:." + str(self.decimals) + "f}").format(100 * (iteration / float(self.total)))
        filled_length = int(self.length * iteration // self.total)
        bar = self.fill * filled_length + '-' * (self.length - filled_length)
        
        # Цветная версия
        bar_colored = f"\033[0;32m{self.fill * filled_length}\033[0;37m{'-' * (self.length - filled_length)}\033[0m"
        
        print(f'\r{self.prefix} |{bar_colored}| {percent}% {self.suffix}', end=self.print_end)
        
        if iteration == self.total:
            print()
    
    def close(self):
        """Завершает прогресс-бар"""
        print()


class ASCIIArt:
    """ASCII арт для красивого интерфейса"""
    
    @staticmethod
    def print_header():
        """Печатает главный ASCII заголовок"""
        print("\033[1;35m")  # Яркий пурпурный
        print("""
╔══════════════════════════════════════════════════════════════╗
║                                                              ║
║  ██████╗ ██╗    ██╗    ████████╗███████╗ ██████╗██╗  ██╗    ║
║  ██╔══██╗██║    ██║    ╚══██╔══╝██╔════╝██╔════╝██║  ██║    ║
║  ██████╔╝██║ █╗ ██║       ██║   █████╗  ██║     ███████║    ║
║  ██╔══██╗██║███╗██║       ██║   ██╔══╝  ██║     ██╔══██║    ║
║  ██║  ██║╚███╔███╔╝       ██║   ███████╗╚██████╗██║  ██║    ║
║  ╚═╝  ╚═╝ ╚══╝╚══╝        ╚═╝   ╚══════╝ ╚═════╝╚═╝  ╚═╝    ║
║                                                              ║
║        🚀 RW TECH PPTX GENERATOR 🚀                         ║
║     💎 Universal PowerPoint Creation System 💎             ║
║                                                              ║
╚══════════════════════════════════════════════════════════════╝
""")
        print("\033[1;36m✨ Профессиональная генерация презентаций с ИИ \033[1;35m🪄\033[0m")
    
    @staticmethod
    def print_divider(symbol="═", length=60, color="\033[1;36m"):
        """Печатает разделительную линию"""
        print(f"{color}{'═' * length}\033[0m")
    
    @staticmethod
    def print_box(text, color="\033[1;36m"):
        """Печатает текст в ASCII рамке"""
        length = len(text)
        border = "─" * (length + 2)
        
        print(f"{color}┌{border}┐\033[0m")
        print(f"{color}│ {text} │\033[0m")
        print(f"{color}└{border}┘\033[0m")
    
    @staticmethod
    def print_loading_frame(frame_num):
        """Печатает кадры анимации загрузки"""
        frames = [
            "⠋", "⠙", "⠹", "⠸", "⠼", "⠴", "⠦", "⠧",
            "⠇", "⠏", "⠋", "⠙", "⠹", "⠸", "⠼", "⠴"
        ]
        return frames[frame_num % len(frames)]
    
    @staticmethod
    def animate_text(text, color="\033[1;36m", delay=0.03):
        """Анимированный вывод текста по символам"""
        import time
        for char in text:
            print(f"{color}{char}\033[0m", end='', flush=True)
            time.sleep(delay)
        print()
    
    @staticmethod
    def print_rw_tech_logo():
        """Печатает красивый логотип RW Tech"""
        print("\033[1;34m")  # Яркий синий
        print("""
    ╔═══════════════════════════════════════════════════════╗
    ║                                                       ║
    ║    ██████╗ ██╗    ██╗    ████████╗███████╗ ██████╗    ║
    ║    ██╔══██╗██║    ██║    ╚══██╔══╝██╔════╝██╔════╝    ║
    ║    ██████╔╝██║ █╗ ██║       ██║   █████╗  ██║         ║
    ║    ██╔══██╗██║███╗██║       ██║   ██╔══╝  ██║         ║
    ║    ██║  ██║╚███╔███╔╝       ██║   ███████╗╚██████╗    ║
    ║    ╚═╝  ╚═╝ ╚══╝╚══╝        ╚═╝   ╚══════╝ ╚═════╝    ║
    ║                                                       ║
    ║              \033[1;36m🌟 REVOLUTIONARY WORKFLOWS 🌟\033[1;34m          ║
    ║              \033[1;35m💎 TECHNOLOGY SOLUTIONS 💎\033[1;34m           ║
    ║                                                       ║
    ╚═══════════════════════════════════════════════════════╝
        """)
        print("\033[0m")
    
    @staticmethod 
    def print_success_banner():
        """Печатает баннер успешного завершения"""
        print("\033[1;32m")  # Яркий зеленый
        print("""
    ╔═════════════════════════════════════════════════════════════╗
    ║                                                             ║
    ║   🎉 PRESENTATION GENERATED SUCCESSFULLY! 🎉               ║
    ║                                                             ║
    ║     ✨ Powered by RW Tech AI Technology ✨                 ║
    ║     🚀 Next-Generation Content Creation 🚀                 ║
    ║                                                             ║
    ║              Thank you for using RW Tech!                   ║
    ║                                                             ║
    ╚═════════════════════════════════════════════════════════════╝
        """)
        print("\033[0m")


class ColorfulUI:
    """Класс для красивого цветного вывода с анимациями"""
    
    # ANSI цветовые коды
    COLORS = {
        'red': '\033[0;31m',
        'green': '\033[0;32m', 
        'yellow': '\033[1;33m',
        'blue': '\033[0;34m',
        'purple': '\033[0;35m',
        'cyan': '\033[0;36m',
        'white': '\033[1;37m',
        'bold': '\033[1m',
        'dim': '\033[2m',
        'reset': '\033[0m',
        'bright_green': '\033[1;32m',
        'bright_blue': '\033[1;34m',
        'bright_yellow': '\033[1;33m',
        'bright_cyan': '\033[1;36m',
        'bright_magenta': '\033[1;35m'
    }
    
    # Эмодзи для разных статусов
    EMOJI = {
        'success': '✅',
        'error': '❌', 
        'warning': '⚠️',
        'info': 'ℹ️',
        'progress': '🔄',
        'fire': '🔥',
        'rocket': '🚀',
        'star': '⭐',
        'gem': '💎',
        'art': '🎨',
        'magic': '✨',
        'lightning': '⚡',
        'crown': '👑',
        'trophy': '🏆',
        'brain': '🧠',
        'paint': '🎨',
        'wand': '🪄',
        'tech': '⚙️',
        'innovative': '🌟',
        'workflow': '🔄',
        'revolution': '💫'
    }
    
    # RW Tech брендинг
    RW_TECH_COLORS = {
        'primary': '\033[1;34m',      # Яркий синий
        'secondary': '\033[1;36m',     # Яркий голубой  
        'accent': '\033[1;35m',        # Яркий фиолетовый
        'success': '\033[1;32m',       # Яркий зеленый
        'warning': '\033[1;33m',       # Яркий желтый
        'danger': '\033[1;31m',        # Яркий красный
        'gradient1': '\033[38;5;33m',  # Градиент синий
        'gradient2': '\033[38;5;39m',  # Градиент голубой
        'gradient3': '\033[38;5;45m'   # Градиент светло-голубой
    }
    
    # Анимированный спиннер
    SPINNER_FRAMES = ["⠋", "⠙", "⠹", "⠸", "⠼", "⠴", "⠦", "⠧"]
    
    @classmethod
    def animate_spinner(cls, text, duration=2):
        """Показывает анимированный спиннер"""
        import time
        frame = 0
        end_time = time.time() + duration
        
        while time.time() < end_time:
            spinner = cls.SPINNER_FRAMES[frame % len(cls.SPINNER_FRAMES)]
            print(f"\r{cls.COLORS['cyan']}{spinner} {text}{cls.COLORS['reset']}", end='', flush=True)
            time.sleep(0.1)
            frame += 1
        
        print(f"\r{cls.COLORS['green']}✅ {text} завершено{cls.COLORS['reset']}")
    
    @classmethod
    def print_ascii_step(cls, step_num, title, description=""):
        """Печатает красивый ASCII этап"""
        print()
        cls.print_divider("▓", 60, cls.COLORS['bright_blue'])
        print(f"{cls.COLORS['bright_cyan']}🔸 ЭТАП {step_num}: {cls.COLORS['bold']}{title}{cls.COLORS['reset']}")
        if description:
            print(f"{cls.COLORS['dim']}   {description}{cls.COLORS['reset']}")
        cls.print_divider("▓", 60, cls.COLORS['bright_blue'])
    
    @classmethod
    def print_divider(cls, symbol="═", length=60, color=None):
        """Печатает разделительную линию"""
        if color is None:
            color = cls.COLORS['cyan']
        print(f"{color}{symbol * length}{cls.COLORS['reset']}")
    
    @classmethod
    def print_rw_tech_banner(cls):
        """Печатает красивый баннер RW Tech с анимацией"""
        import time
        
        # Заголовок с градиентом
        print(f"\n{cls.RW_TECH_COLORS['gradient1']}╔{'═' * 70}╗{cls.COLORS['reset']}")
        print(f"{cls.RW_TECH_COLORS['gradient2']}║{' ' * 70}║{cls.COLORS['reset']}")
        print(f"{cls.RW_TECH_COLORS['gradient3']}║   {cls.EMOJI['rocket']} {cls.RW_TECH_COLORS['primary']}RW TECH - Revolutionary Workflows & Technology{cls.RW_TECH_COLORS['gradient3']}   ║{cls.COLORS['reset']}")
        print(f"{cls.RW_TECH_COLORS['gradient2']}║{' ' * 70}║{cls.COLORS['reset']}")
        print(f"{cls.RW_TECH_COLORS['gradient1']}╚{'═' * 70}╝{cls.COLORS['reset']}")
        
        # Анимированный слоган
        slogans = [
            f"{cls.EMOJI['magic']} Transforming Ideas into Reality",
            f"{cls.EMOJI['lightning']} Next-Generation AI Solutions", 
            f"{cls.EMOJI['innovative']} Innovation at Every Step",
            f"{cls.EMOJI['workflow']} Streamlined Digital Workflows"
        ]
        
        for slogan in slogans:
            print(f"\r{cls.RW_TECH_COLORS['secondary']}{slogan:^70}{cls.COLORS['reset']}", end='', flush=True)
            time.sleep(0.8)
        print()
    
    @classmethod
    def print_rw_tech_step(cls, step_name, description=""):
        """Печатает этап в стиле RW Tech"""
        print(f"\n{cls.RW_TECH_COLORS['primary']}{'▓' * 60}{cls.COLORS['reset']}")
        print(f"{cls.RW_TECH_COLORS['secondary']}{cls.EMOJI['tech']} {step_name}{cls.COLORS['reset']}")
        if description:
            print(f"{cls.COLORS['dim']}{description}{cls.COLORS['reset']}")
        print(f"{cls.RW_TECH_COLORS['primary']}{'▓' * 60}{cls.COLORS['reset']}")
    
    @classmethod
    def print_rw_tech_success(cls, message):
        """Печатает сообщение об успехе в стиле RW Tech"""
        print(f"{cls.RW_TECH_COLORS['success']}{cls.EMOJI['success']} {message}{cls.COLORS['reset']}")
    
    @classmethod
    def print_rw_tech_error(cls, message):
        """Печатает сообщение об ошибке в стиле RW Tech"""
        print(f"{cls.RW_TECH_COLORS['danger']}{cls.EMOJI['error']} {message}{cls.COLORS['reset']}")
    
    @classmethod
    def print_rw_tech_info(cls, message):
        """Печатает информационное сообщение в стиле RW Tech"""
        print(f"{cls.RW_TECH_COLORS['secondary']}{cls.EMOJI['info']} {message}{cls.COLORS['reset']}")
    
    @classmethod
    def animated_rw_tech_loading(cls, text, duration=3):
        """Анимированная загрузка в стиле RW Tech"""
        import time
        frames = ["⚙️", "🔧", "⚡", "🌟", "💫", "✨"]
        frame = 0
        end_time = time.time() + duration
        
        while time.time() < end_time:
            emoji = frames[frame % len(frames)]
            print(f"\r{cls.RW_TECH_COLORS['accent']}{emoji} {text}...{cls.COLORS['reset']}", end='', flush=True)
            time.sleep(0.3)
            frame += 1
        
        print(f"\r{cls.RW_TECH_COLORS['success']}{cls.EMOJI['success']} {text} завершено{cls.COLORS['reset']}")
    
    @classmethod
    def print_progress_animation(cls, text, progress=0, total=100):
        """Печатает анимированный прогресс"""
        bar_length = 40
        filled = int(bar_length * progress / total)
        bar = "█" * filled + "▒" * (bar_length - filled)
        percentage = int(100 * progress / total)
        
        print(f"\r{cls.COLORS['bright_cyan']}{text} [{cls.COLORS['bright_green']}{bar}{cls.COLORS['bright_cyan']}] {percentage}%{cls.COLORS['reset']}", end='', flush=True)
        
        if progress >= total:
            print()  # Новая строка в конце
    
    @classmethod 
    def print_header(cls, text, color='purple', emoji='rocket'):
        """Печатает красивый заголовок"""
        emoji_char = cls.EMOJI.get(emoji, '🚀')
        color_code = cls.COLORS.get(color, cls.COLORS['purple'])
        reset = cls.COLORS['reset']
        
        print(f"\n{color_code}{'='*60}{reset}")
        print(f"{color_code}{emoji_char} {text} {emoji_char}{reset}")
        print(f"{color_code}{'='*60}{reset}")
    
    @classmethod
    def print_step(cls, step_num, title, color='blue', emoji='info'):
        """Печатает этап выполнения"""
        emoji_char = cls.EMOJI.get(emoji, 'ℹ️')
        color_code = cls.COLORS.get(color, cls.COLORS['blue'])
        reset = cls.COLORS['reset']
        
        print(f"\n{color_code}📋 ЭТАП {step_num}: {title}{reset}")
        print(f"{color_code}{'─'*50}{reset}")
    
    @classmethod
    def print_success(cls, text, emoji='success'):
        """Печатает сообщение об успехе"""
        emoji_char = cls.EMOJI.get(emoji, '✅')
        color = cls.COLORS['green']
        reset = cls.COLORS['reset']
        print(f"{color}{emoji_char} {text}{reset}")
    
    @classmethod
    def print_error(cls, text, emoji='error'):
        """Печатает сообщение об ошибке"""
        emoji_char = cls.EMOJI.get(emoji, '❌')
        color = cls.COLORS['red']
        reset = cls.COLORS['reset']
        print(f"{color}{emoji_char} {text}{reset}")
    
    @classmethod
    def print_warning(cls, text, emoji='warning'):
        """Печатает предупреждение"""
        emoji_char = cls.EMOJI.get(emoji, '⚠️')
        color = cls.COLORS['yellow']
        reset = cls.COLORS['reset']
        print(f"{color}{emoji_char} {text}{reset}")
    
    @classmethod 
    def print_prompt_generation(cls, slide_num, total):
        """Печатает сообщение о генерации промптов"""
        color = cls.COLORS['bright_cyan']
        reset = cls.COLORS['reset']
        emoji = cls.EMOJI['brain']
        print(f"{color}{emoji} Генерация промпта для слайда {slide_num}/{total}...{reset}")
    
    @classmethod
    def print_image_generation(cls, slide_num, total, model_name):
        """Печатает сообщение о генерации изображений"""
        color = cls.COLORS['bright_magenta']
        reset = cls.COLORS['reset']
        emoji = cls.EMOJI['wand']
        print(f"{color}{emoji} Генерация изображения {slide_num}/{total} с помощью {model_name}...{reset}")
    
    @classmethod
    def create_progress_bar(cls, total, prefix, color='green'):
        """Создает цветной прогресс-бар"""
        return ProgressBar(total, prefix=f"{cls.COLORS[color]}{prefix}{cls.COLORS['reset']}")
    
    @classmethod
    def print_info(cls, text, emoji='info', color='cyan'):
        """Печатает информационное сообщение"""
        emoji_char = cls.EMOJI.get(emoji, 'ℹ️')
        color_code = cls.COLORS.get(color, cls.COLORS['cyan'])
        reset = cls.COLORS['reset']
        print(f"{color_code}{emoji_char} {text}{reset}")
    
    @classmethod
    def print_progress(cls, text, emoji='progress'):
        """Печатает прогресс с анимацией"""
        emoji_char = cls.EMOJI.get(emoji, '🔄')
        color = cls.COLORS['blue']
        reset = cls.COLORS['reset']
        print(f"{color}{emoji_char} {text}...{reset}")
    
    @classmethod
    def animate_loading(cls, text, duration=1.0):
        """Анимация загрузки"""
        frames = ['⠋', '⠙', '⠹', '⠸', '⠼', '⠴', '⠦', '⠧', '⠇', '⠏']
        color = cls.COLORS['cyan']
        reset = cls.COLORS['reset']
        
        start_time = time.time()
        i = 0
        while time.time() - start_time < duration:
            frame = frames[i % len(frames)]
            print(f"\r{color}{frame} {text}...{reset}", end='', flush=True)
            time.sleep(0.1)
            i += 1
        print(f"\r{cls.COLORS['green']}✅ {text} завершено{reset}")
    
    @classmethod
    def print_model_choice(cls, number, name, description, emoji='star'):
        """Печатает вариант выбора модели"""
        emoji_char = cls.EMOJI.get(emoji, '⭐')
        bold = cls.COLORS['bold']
        cyan = cls.COLORS['cyan']
        reset = cls.COLORS['reset']
        
        print(f"{cyan}{number}. {emoji_char} {bold}{name}{reset}{cyan} - {description}{reset}")
        
    @classmethod
    def print_stats(cls, stats_dict):
        """Печатает красивую статистику"""
        print(f"\n{cls.COLORS['purple']}📊 СТАТИСТИКА ВЫПОЛНЕНИЯ{cls.COLORS['reset']}")
        print(f"{cls.COLORS['purple']}{'─'*40}{cls.COLORS['reset']}")
        
        for key, value in stats_dict.items():
            emoji_char = '📈' if isinstance(value, (int, float)) and value > 0 else '📊'
            print(f"{cls.COLORS['cyan']}{emoji_char} {key}: {cls.COLORS['white']}{value}{cls.COLORS['reset']}")


class PromptTemplates:
    """Шаблоны промптов для разных типов визуализации"""
    
    # Обязательная инструкция для всех промптов
    MANDATORY_INSTRUCTION = "Image on WHITE BACKGROUND. Photorealistic quality, hi-res 4K resolution, content highly relevant to context and instructions."
    
    TECHNICAL_ARCHITECTURE = """
    {concept} ARCHITECTURE DIAGRAM, isometric 3D view, WHITE BACKGROUND, 
    connected components with arrows, tech stack visualization, 
    blue and gray color scheme, ultra minimalist design, 
    studio lighting, professional technical illustration,
    Figma style, clean lines, HIGH DETAIL. """ + MANDATORY_INSTRUCTION
    
    PROCESS_FLOW = """
    {process} FLOWCHART, step-by-step visualization, WHITE BACKGROUND,
    numbered stages, arrows showing flow direction, 
    gradient colors from start to finish, modern flat design,
    professional business diagram, ultra high definition,
    clean geometric shapes, CLEAR LABELS. """ + MANDATORY_INSTRUCTION
    
    DATA_VISUALIZATION = """
    {data_concept} INFOGRAPHIC, data visualization, WHITE BACKGROUND,
    modern charts and graphs, bright accent colors on white,
    minimalist style, professional statistics presentation,
    3D isometric elements, studio lighting, ULTRA CLEAN design,
    business analytics style. """ + MANDATORY_INSTRUCTION
    
    CONCEPT_METAPHOR = """
    {concept} visual METAPHOR, symbolic representation, WHITE BACKGROUND,
    minimalist conceptual art, professional illustration,
    single focal point, bright accent color, ultra simple design,
    studio lighting, HIGH QUALITY rendering, 
    modern business presentation style. """ + MANDATORY_INSTRUCTION
    
    ICON_SET = """
    {elements} ICON SET, flat design icons, WHITE BACKGROUND,
    consistent style, bright gradient colors, 
    rounded corners, professional UI elements,
    grid layout, ultra minimalist, studio lighting,
    dribble style, HIGH RESOLUTION. """ + MANDATORY_INSTRUCTION
    
    COMPARISON = """
    {item1} VS {item2} COMPARISON, split screen visualization, WHITE BACKGROUND,
    side by side elements, contrasting colors,
    professional infographic style, clean design,
    isometric 3D view, studio lighting, CLEAR DIFFERENCES highlighted. """ + MANDATORY_INSTRUCTION
    
    TIMELINE = """
    {topic} TIMELINE visualization, chronological progression, WHITE BACKGROUND,
    horizontal or vertical timeline, milestone markers,
    gradient color progression, modern design,
    professional presentation style, ultra clean,
    HIGH DETAIL, clear date labels. """ + MANDATORY_INSTRUCTION


class ExecutionStats:
    """Класс для подробной статистики выполнения"""
    
    def __init__(self):
        self.stats = {
            'start_time': datetime.now(),
            'api_validation_attempts': 0,
            'api_validation_success': False,
            'slides_to_process': 0,
            'prompts_attempted': 0,
            'prompts_generated': 0,
            'prompts_failed': 0,
            'images_attempted': 0,
            'images_generated': 0,
            'images_failed': 0,
            'images_inserted': 0,
            'total_api_calls': 0,
            'total_errors': 0,
            'end_time': None,
            'total_duration': None
        }
    
    def increment(self, stat_name, amount=1):
        """Увеличивает значение статистики"""
        if stat_name in self.stats:
            self.stats[stat_name] += amount
    
    def set(self, stat_name, value):
        """Устанавливает значение статистики"""
        self.stats[stat_name] = value
    
    def get(self, stat_name, default=0):
        """Получает значение статистики"""
        return self.stats.get(stat_name, default)
    
    def _calculate_success_rate(self, operation_type):
        """Вычисляет коэффициент успешности"""
        if operation_type == 'prompts':
            total = self.stats['prompts_attempted']
            success = self.stats['prompts_generated']
        elif operation_type == 'images':
            total = self.stats['images_attempted']
            success = self.stats['images_generated']
        else:
            return 0.0
        
        return (success / total) if total > 0 else 0.0
    
    def _is_execution_successful(self):
        """Определяет, было ли выполнение успешным"""
        prompts_rate = self._calculate_success_rate('prompts')
        images_rate = self._calculate_success_rate('images')
        
        # Требуем минимум 80% успешности
        return (self.stats['api_validation_success'] and 
                prompts_rate >= 0.8 and 
                images_rate >= 0.8 and
                self.stats['images_inserted'] > 0)
    
    def print_final_report(self):
        """Выводит детальный финальный отчет"""
        self.stats['end_time'] = datetime.now()
        self.stats['total_duration'] = self.stats['end_time'] - self.stats['start_time']
        
        print("\n" + "="*60)
        print("📊 ФИНАЛЬНЫЙ ОТЧЕТ ВЫПОЛНЕНИЯ")
        print("="*60)
        
        print(f"\n⏱️  Время выполнения: {self.stats['total_duration']}")
        print(f"\n🔑 Валидация API:")
        print(f"   Попыток: {self.stats['api_validation_attempts']}")
        print(f"   Результат: {'✅ Успешно' if self.stats['api_validation_success'] else '❌ Неудачно'}")
        
        if self.stats['slides_to_process'] > 0:
            print(f"\n📝 Генерация промптов:")
            print(f"   Запланировано: {self.stats['slides_to_process']}")
            print(f"   Попыток: {self.stats['prompts_attempted']}")
            print(f"   Успешно: {self.stats['prompts_generated']}")
            print(f"   Неудачно: {self.stats['prompts_failed']}")
            print(f"   Успешность: {self._calculate_success_rate('prompts'):.1%}")
            
            print(f"\n🎨 Генерация изображений:")
            print(f"   Попыток: {self.stats['images_attempted']}")
            print(f"   Успешно: {self.stats['images_generated']}")
            print(f"   Неудачно: {self.stats['images_failed']}")
            print(f"   Вставлено в презентацию: {self.stats['images_inserted']}")
            print(f"   Успешность: {self._calculate_success_rate('images'):.1%}")
        
        print(f"\n📊 Общая статистика:")
        print(f"   API вызовов: {self.stats['total_api_calls']}")
        print(f"   Ошибок: {self.stats['total_errors']}")
        
        # Итоговый вердикт
        if self._is_execution_successful():
            print("\n✅ ВЫПОЛНЕНИЕ УСПЕШНО ЗАВЕРШЕНО")
            return True
        else:
            print("\n❌ ВЫПОЛНЕНИЕ ЗАВЕРШЕНО С ОШИБКАМИ")
            return False


class ExecutionCheckpoints:
    """Класс для контроля выполнения через контрольные точки"""
    
    def __init__(self, stats, parent=None):
        self.stats = stats
        self.parent = parent
        self.checkpoints = {
            'api_validation': False,
            'prompts_generation': False,
            'images_generation': False,
            'presentation_update': False,
            'final_validation': False
        }
    
    def validate_checkpoint(self, checkpoint_name, required_success_rate=0.8):
        """Валидирует контрольную точку"""
        print(f"\n🔍 Проверка контрольной точки: {checkpoint_name}")
        
        if checkpoint_name == 'api_validation':
            success = self.stats.get('api_validation_success')
            if not success:
                print(f"⛔ КРИТИЧЕСКАЯ ОШИБКА: API ключи не прошли валидацию")
                return False
        
        elif checkpoint_name == 'prompts_generation':
            total = self.stats.get('prompts_attempted')
            success = self.stats.get('prompts_generated')
            success_rate = (success / total) if total > 0 else 0.0
            
            if success_rate < required_success_rate:
                print(f"⛔ КРИТИЧЕСКАЯ ОШИБКА на этапе '{checkpoint_name}'")
                print(f"   Успешность: {success_rate:.1%} (требуется минимум {required_success_rate:.1%})")
                print(f"   Успешно: {success}")
                print(f"   Неудачно: {self.stats.get('prompts_failed')}")
                print(f"   Из {total} попыток")
                return False
        
        elif checkpoint_name == 'images_generation':
            total = self.stats.get('images_attempted')
            success = self.stats.get('images_generated')
            success_rate = (success / total) if total > 0 else 0.0
            
            if success_rate < required_success_rate:
                print(f"⛔ КРИТИЧЕСКАЯ ОШИБКА на этапе '{checkpoint_name}'")
                print(f"   Успешность: {success_rate:.1%} (требуется минимум {required_success_rate:.1%})")
                print(f"   Успешно: {success}")
                print(f"   Неудачно: {self.stats.get('images_failed')}")
                print(f"   Из {total} попыток")
                return False
        
        elif checkpoint_name == 'presentation_update':
            inserted = self.stats.get('images_inserted')
            generated = self.stats.get('images_generated')
            
            # Получаем путь к result_dir от родительского объекта
            if self.parent and hasattr(self.parent, 'result_dir'):
                illustrated_path = os.path.join(self.parent.result_dir, "RWTech_Universal_Presentation_Illustrated.pptx")
            else:
                illustrated_path = os.path.join(os.getcwd(), "pptx_result", "RWTech_Universal_Presentation_Illustrated.pptx")
            
            if generated > 0:
                if inserted == 0:
                    print(f"⛔ КРИТИЧЕСКАЯ ОШИБКА: Изображения созданы ({generated}), но не вставлены в презентацию")
                    return False
                elif not os.path.exists(illustrated_path):
                    print(f"⛔ КРИТИЧЕСКАЯ ОШИБКА: Файл с иллюстрациями не создан: {illustrated_path}")
                    return False
                else:
                    # Проверяем что презентация действительно содержит изображения
                    try:
                        from pptx import Presentation
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        
                        prs = Presentation(illustrated_path)
                        images_found = 0
                        
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    images_found += 1
                        
                        # Проверяем количество изображений (должно быть больше шаблонных)
                        template_images_count = len(prs.slides) * 2  # 2 изображения на слайд из шаблона
                        ai_images_found = images_found - template_images_count
                        
                        if ai_images_found < inserted:
                            print(f"⛔ КРИТИЧЕСКАЯ ОШИБКА: В презентации найдено {ai_images_found} AI-изображений, ожидалось {inserted}")
                            return False
                        else:
                            print(f"✅ Проверка презентации: найдено {ai_images_found} AI-изображений из {inserted} вставленных")
                            
                    except Exception as e:
                        print(f"⛔ ОШИБКА проверки презентации: {e}")
                        return False
        
        self.checkpoints[checkpoint_name] = True
        print(f"✅ Контрольная точка '{checkpoint_name}' пройдена успешно")
        return True
    
    def get_checkpoint_status(self):
        """Возвращает статус всех контрольных точек"""
        return self.checkpoints.copy()


class RWTechPPTXGenerator:
    def __init__(self):
        self.base_path = os.getcwd()  # Use current working directory
        self.content_file = os.path.join(self.base_path, "pptx_content", "slide_content.txt")
        self.template_file = os.path.join(self.base_path, "pptx_template", "Шаблон презентации 16х9.pptx")
        self.result_dir = os.path.join(self.base_path, "pptx_result")
        self.result_file = os.path.join(self.result_dir, "RWTech_Universal_Presentation.pptx")
        
        # AI illustration directories
        self.prompts_dir = os.path.join(self.base_path, "prompts_for_img")
        self.images_dir = os.path.join(self.base_path, "img_generated")
        self.logs_dir = os.path.join(self.base_path, "logs")
        
        # Initialize logging
        self.logger = None
        self._setup_logging()
        
        # API credentials (will be loaded from config)
        self.claude_api_key = ""
        self.openai_api_key = ""
        self.gemini_api_key = ""
        self._load_config()
        
        # AI settings
        self.use_ai_illustrations = False
        self.slide_interval = 5  # Every 5th slide by default
        self.image_model = 'dall-e-3'  # Default image generation model
        
        # New execution control systems
        self.execution_stats = ExecutionStats()
        self.checkpoints = ExecutionCheckpoints(self.execution_stats, self)
        
        # Legacy stats for backward compatibility
        self.generation_stats = {
            'total_slides': 0,
            'prompts_generated': 0,
            'prompts_failed': 0,
            'images_generated': 0,
            'images_failed': 0
        }
        
        self.slides_data = []
        self.template_images = []
        self.prs = None  # Для хранения ссылки на презентацию
        
    def validate_files(self):
        """Проверка существования всех необходимых файлов"""
        print("Проверка файлов...")
        
        if not os.path.exists(self.content_file):
            print(f"ОШИБКА: Файл с контентом не найден: {self.content_file}")
            sys.exit(1)
            
        if not os.path.exists(self.template_file):
            print(f"ОШИБКА: Файл шаблона не найден: {self.template_file}")
            sys.exit(1)
            
        # Создаем директорию результата если не существует
        if not os.path.exists(self.result_dir):
            os.makedirs(self.result_dir)
            print(f"Создана директория: {self.result_dir}")
        
        # Создаем директории для AI иллюстраций
        if not os.path.exists(self.prompts_dir):
            os.makedirs(self.prompts_dir)
            print(f"Создана директория для промптов: {self.prompts_dir}")
            
        if not os.path.exists(self.images_dir):
            os.makedirs(self.images_dir)
            print(f"Создана директория для изображений: {self.images_dir}")
            
        if not os.path.exists(self.logs_dir):
            os.makedirs(self.logs_dir)
            print(f"Создана директория для логов: {self.logs_dir}")
            
        print("Все файлы найдены ✓")

    def setup_ai_illustrations(self, interactive=True):
        """Интерактивная настройка AI иллюстраций"""
        print("\n=== Настройка AI-иллюстраций ===")
        
        if not interactive:
            print("Режим без интерактивного ввода: AI-иллюстрации отключены")
            self.use_ai_illustrations = False
            return
            
        print("Хотите ли вы генерировать AI-иллюстрации для слайдов?")
        print("Это добавит визуальные элементы к презентации с помощью DALL-E 3.")
        
        try:
            while True:
                choice = input("Генерировать AI-иллюстрации? (да/нет): ").strip().lower()
                if choice in ['да', 'yes', 'y', '1']:
                    self.use_ai_illustrations = True
                    break
                elif choice in ['нет', 'no', 'n', '0']:
                    self.use_ai_illustrations = False
                    print("AI-иллюстрации отключены. Будет использован стандартный режим.")
                    return
                else:
                    print("Пожалуйста, введите 'да' или 'нет'")
        except EOFError:
            print("Автоматический режим: AI-иллюстрации отключены")
            self.use_ai_illustrations = False
            return
        
        if self.use_ai_illustrations:
            print("\nВыберите режим генерации иллюстраций:")
            print("1. Каждый 3-й слайд")
            print("2. Каждый 5-й слайд (рекомендуется)")
            print("3. Каждый 10-й слайд")
            print("4. Пользовательский интервал")
            
            try:
                while True:
                    choice = input("Ваш выбор (1-4): ").strip()
                    if choice == '1':
                        self.slide_interval = 3
                        break
                    elif choice == '2':
                        self.slide_interval = 5
                        break
                    elif choice == '3':
                        self.slide_interval = 10
                        break
                    elif choice == '4':
                        while True:
                            try:
                                custom_interval = int(input("Введите интервал (например, 7 для каждого 7-го слайда): "))
                                if custom_interval > 0:
                                    self.slide_interval = custom_interval
                                    break
                                else:
                                    print("Интервал должен быть положительным числом.")
                            except (ValueError, EOFError):
                                print("Использован интервал по умолчанию: 5")
                                self.slide_interval = 5
                                break
                        break
                    else:
                        print("Пожалуйста, введите число от 1 до 4")
            except EOFError:
                print("Использован интервал по умолчанию: каждый 5-й слайд")
                self.slide_interval = 5
            
            print(f"\n✓ AI-иллюстрации будут генерироваться для каждого {self.slide_interval}-го слайда")
            
            # Выбор модели для генерации изображений
            ColorfulUI.print_header("Выбор AI-модели для генерации изображений", emoji='art')
            
            ColorfulUI.print_model_choice("1", "DALL-E 3 (рекомендуется)", 
                                        "высокое качество, отличное понимание промптов", emoji='fire')
            ColorfulUI.print_model_choice("2", "GPT-Image-1 (новая)", 
                                        "быстрая генерация, больше возможностей настройки", emoji='lightning')  
            ColorfulUI.print_model_choice("3", "Google Gemini 2.0 Flash (контекстные)", 
                                        "контекстно-релевантные изображения с рассуждениями", emoji='magic')
            ColorfulUI.print_model_choice("4", "Google Imagen 3 (художественные)", 
                                        "фотореализм, художественная детализация", emoji='gem')
            
            try:
                while True:
                    model_choice = input(f"\n{ColorfulUI.COLORS['cyan']}🎯 Ваш выбор (1-4): {ColorfulUI.COLORS['reset']}").strip()
                    if model_choice == '1':
                        self.image_model = 'dall-e-3'
                        ColorfulUI.print_success("Выбрана модель DALL-E 3", emoji='fire')
                        break
                    elif model_choice == '2':
                        self.image_model = 'gpt-image-1'
                        ColorfulUI.print_success("Выбрана модель GPT-Image-1", emoji='lightning')
                        break
                    elif model_choice == '3':
                        self.image_model = 'gemini-2.0-flash'
                        ColorfulUI.print_success("Выбрана модель Google Gemini 2.0 Flash", emoji='magic')
                        break
                    elif model_choice == '4':
                        self.image_model = 'imagen-3'
                        ColorfulUI.print_success("Выбрана модель Google Imagen 3", emoji='gem')
                        break
                    else:
                        ColorfulUI.print_error("Пожалуйста, введите 1, 2, 3 или 4")
            except EOFError:
                ColorfulUI.print_warning("Использована модель по умолчанию: DALL-E 3")
                self.image_model = 'dall-e-3'
            
            # Сохраняем выбранную модель в конфигурацию
            self._save_config()
            
            print("✓ Специальные слайды (титульный, цитаты, перерывы) будут пропущены")
            print("✓ Изображения будут сохранены в директории img_generated/")
            print("✓ Промпты для генерации будут сохранены в директории prompts_for_img/")

    def _setup_logging(self):
        """Настраивает систему логирования"""
        try:
            if not hasattr(self, 'logs_dir'):
                self.logs_dir = os.path.join(self.base_path, "logs")
            
            os.makedirs(self.logs_dir, exist_ok=True)
            
            # Имя файла лога с датой
            log_filename = f"generation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"
            log_path = os.path.join(self.logs_dir, log_filename)
            
            # Настройка логгера
            logging.basicConfig(
                level=logging.INFO,
                format='%(asctime)s - %(levelname)s - %(message)s',
                handlers=[
                    logging.FileHandler(log_path, encoding='utf-8'),
                    logging.StreamHandler()
                ]
            )
            
            self.logger = logging.getLogger(__name__)
            self.logger.info(f"Логирование настроено. Файл лога: {log_path}")
            
        except Exception as e:
            print(f"Предупреждение: Не удалось настроить логирование: {e}")
            self.logger = logging.getLogger(__name__)

    def _load_config(self):
        """Загружает конфигурацию из файла или переменных окружения"""
        config_path = os.path.join(self.base_path, 'config.json')
        
        # Попытка загрузить из файла
        if os.path.exists(config_path):
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    self.claude_api_key = config.get('claude_api_key', '')
                    self.openai_api_key = config.get('openai_api_key', '')
                    self.gemini_api_key = config.get('gemini_api_key', '')
                    self.image_model = config.get('image_model', 'dall-e-3')
                    if self.logger:
                        self.logger.info("Конфигурация загружена из файла config.json")
            except Exception as e:
                print(f"Предупреждение: Ошибка чтения config.json: {e}")
        
        # Переменные окружения имеют приоритет
        self.claude_api_key = os.environ.get('CLAUDE_API_KEY', self.claude_api_key)
        self.openai_api_key = os.environ.get('OPENAI_API_KEY', self.openai_api_key)
        self.gemini_api_key = os.environ.get('GEMINI_API_KEY', self.gemini_api_key)
        
        # Если ключей нет, используем встроенные значения как fallback
        # API keys should be provided through config.json or environment variables
        # No hardcoded keys for security reasons
        if not self.claude_api_key:
            self.claude_api_key = None  # Will be prompted interactively
        
        if not self.openai_api_key:
            self.openai_api_key = None  # Will be prompted interactively

    def _save_config(self):
        """Сохраняет конфигурацию в файл"""
        config_path = os.path.join(self.base_path, 'config.json')
        config = {
            'claude_api_key': self.claude_api_key,
            'openai_api_key': self.openai_api_key,
            'gemini_api_key': self.gemini_api_key,
            'image_model': self.image_model
        }
        
        try:
            with open(config_path, 'w', encoding='utf-8') as f:
                json.dump(config, f, indent=2)
            print(f"✓ Конфигурация сохранена в {config_path}")
            if self.logger:
                self.logger.info(f"Конфигурация сохранена в {config_path}")
        except Exception as e:
            print(f"Ошибка сохранения конфигурации: {e}")

    def _test_claude_connection(self):
        """Тестирует соединение с Claude API"""
        try:
            url = "https://api.anthropic.com/v1/messages"
            headers = {
                "x-api-key": self.claude_api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json"
            }
            
            # Минимальный тестовый запрос
            data = {
                "model": "claude-3-5-sonnet-20241022",
                "max_tokens": 10,
                "messages": [{"role": "user", "content": "Hi"}]
            }
            
            response = requests.post(url, headers=headers, json=data, timeout=10)
            
            return {
                'success': response.status_code == 200,
                'status_code': response.status_code,
                'error': response.text if response.status_code != 200 else None,
                'endpoint': url
            }
            
        except Exception as e:
            return {
                'success': False,
                'status_code': None,
                'error': str(e),
                'endpoint': url if 'url' in locals() else 'N/A'
            }

    def _test_openai_connection(self):
        """Тестирует соединение с OpenAI API"""
        try:
            url = "https://api.openai.com/v1/models"
            headers = {
                "Authorization": f"Bearer {self.openai_api_key}",
                "Content-Type": "application/json"
            }
            
            response = requests.get(url, headers=headers, timeout=10)
            
            return {
                'success': response.status_code == 200,
                'status_code': response.status_code,
                'error': response.text if response.status_code != 200 else None,
                'endpoint': url
            }
            
        except Exception as e:
            return {
                'success': False,
                'status_code': None,
                'error': str(e),
                'endpoint': url if 'url' in locals() else 'N/A'
            }

    def _validate_and_update_api_keys(self, interactive=True):
        """СТРОГАЯ валидация API ключей с циклом до успеха"""
        if not interactive:
            print("Режим без интерактивного ввода: пропуск валидации API ключей")
            self.execution_stats.set('api_validation_success', False)
            return False
        
        print("\n🔒 СТРОГАЯ ВАЛИДАЦИЯ API КЛЮЧЕЙ")
        print("Программа не будет продолжена без рабочих ключей!")
        
        # Флаги валидности
        claude_valid = False
        openai_valid = False
        gemini_valid = False
        max_attempts = 3
        
        # Определяем, какие API нужно проверить
        need_gemini = (hasattr(self, 'image_model') and self.image_model in ['gemini-2.0-flash', 'imagen-3'])
        
        # Цикл валидации Claude
        claude_attempts = 0
        while not claude_valid and claude_attempts < max_attempts:
            claude_attempts += 1
            self.execution_stats.increment('api_validation_attempts')
            
            if not self.claude_api_key:
                print("\n❌ API ключ Claude отсутствует")
                self.claude_api_key = self._request_claude_key()
                if not self.claude_api_key:
                    continue
            
            print(f"\n🔍 Проверка Claude API (попытка {claude_attempts}/{max_attempts})...")
            test_result = self._test_claude_connection()
            
            if test_result['success']:
                print("✅ Claude API: подключение успешно")
                claude_valid = True
            else:
                error_action = self._handle_api_error('Claude', test_result)
                
                if error_action == 'new_key':
                    self.claude_api_key = self._request_claude_key()
                elif error_action == 'retry':
                    continue
                elif error_action == 'abort':
                    print("\n⛔ Валидация API отменена пользователем")
                    self.execution_stats.set('api_validation_success', False)
                    return False
                else:  # skip
                    break
        
        if not claude_valid:
            print(f"\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Не удалось подключиться к Claude API за {max_attempts} попыток")
            print("Без рабочего ключа Claude генерация AI-иллюстраций невозможна.")
            self.execution_stats.set('api_validation_success', False)
            return False
        
        # Цикл валидации OpenAI (если нужен)
        need_openai = (hasattr(self, 'image_model') and self.image_model in ['dall-e-3', 'gpt-image-1'])
        
        if need_openai:
            openai_attempts = 0
            while not openai_valid and openai_attempts < max_attempts:
                openai_attempts += 1
                self.execution_stats.increment('api_validation_attempts')
                
                if not self.openai_api_key:
                    print("\n❌ API ключ OpenAI отсутствует")
                    self.openai_api_key = self._request_openai_key()
                    if not self.openai_api_key:
                        continue
                
                print(f"\n🔍 Проверка OpenAI API (попытка {openai_attempts}/{max_attempts})...")
                test_result = self._test_openai_connection()
                
                if test_result['success']:
                    print("✅ OpenAI API: подключение успешно")
                    openai_valid = True
                else:
                    error_action = self._handle_api_error('OpenAI', test_result)
                    
                    if error_action == 'new_key':
                        self.openai_api_key = self._request_openai_key()
                    elif error_action == 'retry':
                        continue
                    elif error_action == 'abort':
                        print("\n⛔ Валидация API отменена пользователем")
                        self.execution_stats.set('api_validation_success', False)
                        return False
                    else:  # skip
                        break
            
            if not openai_valid:
                print(f"\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Не удалось подключиться к OpenAI API за {max_attempts} попыток")
                print("Без рабочего ключа OpenAI генерация изображений невозможна.")
                self.execution_stats.set('api_validation_success', False)
                return False
        else:
            openai_valid = True  # Не нужен OpenAI
        
        # Цикл валидации Gemini (если нужен)
        if need_gemini:
            gemini_attempts = 0
            while not gemini_valid and gemini_attempts < max_attempts:
                gemini_attempts += 1
                self.execution_stats.increment('api_validation_attempts')
                
                if not self.gemini_api_key:
                    model_name = "Gemini 2.0 Flash" if self.image_model == 'gemini-2.0-flash' else "Imagen 3"
                    print(f"\n❌ API ключ Google {model_name} отсутствует")
                    self.gemini_api_key = self._request_gemini_key()
                    if not self.gemini_api_key:
                        continue
                
                model_name = "Gemini 2.0 Flash" if self.image_model == 'gemini-2.0-flash' else "Imagen 3"
                print(f"\n🔍 Проверка Google {model_name} API (попытка {gemini_attempts}/{max_attempts})...")
                test_result = self._test_gemini_connection()
                
                if test_result['success']:
                    print(f"✅ Google {model_name} API: подключение успешно")
                    gemini_valid = True
                else:
                    error_action = self._handle_api_error(f'Google {model_name}', test_result)
                    
                    if error_action == 'new_key':
                        self.gemini_api_key = self._request_gemini_key()
                    elif error_action == 'retry':
                        continue
                    elif error_action == 'abort':
                        print("\n⛔ Валидация API отменена пользователем")
                        self.execution_stats.set('api_validation_success', False)
                        return False
                    else:  # skip
                        break
            
            if not gemini_valid:
                model_name = "Gemini 2.0 Flash" if self.image_model == 'gemini-2.0-flash' else "Imagen 3"
                print(f"\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Не удалось подключиться к Google {model_name} API за {max_attempts} попыток")
                print(f"Без рабочего ключа {model_name} генерация изображений невозможна.")
                self.execution_stats.set('api_validation_success', False)
                return False
        else:
            gemini_valid = True  # Не нужен Gemini
        
        # ОБА ключа валидны
        print("\n🎉 ВСЕ API КЛЮЧИ ПРОШЛИ ВАЛИДАЦИЮ!")
        self.execution_stats.set('api_validation_success', True)
        
        # Автоматически сохраняем конфигурацию
        print("✓ Автоматическое сохранение рабочих API ключей...")
        self._save_config()
        
        return True
    
    def _request_claude_key(self):
        """Запрашивает новый ключ Claude у пользователя"""
        try:
            print("\n🔑 Требуется действительный API ключ Claude")
            print("Формат: sk-ant-api03-...")
            key = input("Введите ключ Claude: ").strip()
            if key and key.startswith('sk-ant-'):
                return key
            else:
                print("❌ Неправильный формат ключа Claude")
                return ""
        except EOFError:
            return ""
    
    def _request_openai_key(self):
        """Запрашивает новый ключ OpenAI у пользователя"""
        try:
            print("\n🔑 Требуется действительный API ключ OpenAI")
            print("Формат: sk-proj-... или sk-...")
            key = input("Введите ключ OpenAI: ").strip()
            if key and key.startswith('sk-'):
                return key
            else:
                print("❌ Неправильный формат ключа OpenAI")
                return ""
        except EOFError:
            return ""
    
    def _test_gemini_connection(self):
        """Тестирует соединение с Google Gemini API"""
        try:
            from google import genai
            
            # Инициализируем клиент
            client = genai.Client(api_key=self.gemini_api_key)
            
            # Пробуем получить список моделей
            models = client.models.list()
            
            # Проверяем, что список не пустой
            if models:
                return {
                    'success': True,
                    'status_code': 200,
                    'error': None,
                    'endpoint': 'Google Gemini API'
                }
            else:
                return {
                    'success': False,
                    'status_code': None,
                    'error': 'Пустой список моделей',
                    'endpoint': 'Google Gemini API'
                }
                
        except ImportError:
            return {
                'success': False,
                'status_code': None,
                'error': 'Библиотека google-genai не установлена',
                'endpoint': 'Google Gemini API'
            }
        except Exception as e:
            return {
                'success': False,
                'status_code': None,
                'error': str(e),
                'endpoint': 'Google Gemini API'
            }
    
    def _request_gemini_key(self):
        """Запрашивает новый ключ Google Gemini у пользователя"""
        try:
            print("\n🔑 Требуется действительный API ключ Google Gemini")
            print("Поддерживает: Gemini 2.0 Flash + Imagen 3")
            print("Формат: AIza...")
            print("Получить ключ: https://aistudio.google.com/apikey")
            key = input("Введите ключ Gemini: ").strip()
            if key and key.startswith('AIza'):
                return key
            else:
                print("❌ Неправильный формат ключа Gemini (должен начинаться с 'AIza')")
                return ""
        except EOFError:
            return ""
    
    def _handle_api_error(self, api_name, error_details):
        """Интерактивная обработка ошибок API"""
        print(f"\n❌ Ошибка {api_name} API:")
        print(f"   Код: {error_details.get('status_code', 'N/A')}")
        if error_details.get('error'):
            print(f"   Сообщение: {error_details['error'][:200]}...")
        
        # Специфичные рекомендации
        status_code = error_details.get('status_code')
        if status_code == 401:
            print("\n📌 Проблема с авторизацией. Возможные причины:")
            print("   • Недействительный API ключ")
            print("   • Ключ был отозван или деактивирован")
            print("   • Неправильный формат ключа")
            print("   • Перепутаны ключи (Claude вместо OpenAI или наоборот)")
        elif status_code == 429:
            print("\n📌 Превышен лимит запросов. Возможные решения:")
            print("   • Подождать несколько минут")
            print("   • Проверить баланс счета")
            print("   • Использовать другой API ключ")
        elif status_code == 404:
            print("\n📌 Неправильный endpoint или отсутствие доступа:")
            print("   • Проверьте правильность API ключа")
            print("   • Убедитесь, что ключ имеет нужные разрешения")
        else:
            print(f"\n📌 Неожиданная ошибка: {status_code}")
        
        # Опции действий
        print(f"\nЧто делать дальше?")
        print("1. Ввести новый API ключ")
        print("2. Повторить попытку с текущим ключом")
        print("3. Отменить выполнение программы")
        
        try:
            choice = input("\nВаш выбор (1-3): ").strip()
            
            if choice == '1':
                return 'new_key'
            elif choice == '2':
                return 'retry'
            else:
                return 'abort'
        except EOFError:
            return 'abort'

    def _log_error(self, error_type, details):
        """Записывает ошибку в лог"""
        try:
            if self.logger:
                self.logger.error(f"{error_type}: {json.dumps(details, ensure_ascii=False, indent=2)}")
        except Exception:
            pass

    def parse_content(self):
        """Парсинг содержимого slide_content.txt"""
        print("Парсинг контента...")
        
        try:
            with open(self.content_file, 'r', encoding='utf-8') as f:
                content = f.read()
        except Exception as e:
            print(f"ОШИБКА при чтении файла контента: {e}")
            sys.exit(1)
            
        # Регулярное выражение для поиска слайдов
        slide_pattern = r'### СЛАЙД (\d+): (.+?)(?=### СЛАЙД \d+:|$)'
        slides = re.findall(slide_pattern, content, re.DOTALL)
        
        if not slides:
            print("ОШИБКА: Не найдено ни одного слайда в файле контента")
            sys.exit(1)
            
        for slide_num, slide_content in slides:
            slide_data = self._parse_slide_content(int(slide_num), slide_content.strip())
            if slide_data:
                self.slides_data.append(slide_data)
                
        print(f"Найдено слайдов: {len(self.slides_data)}")

    def _parse_slide_content(self, slide_num, content):
        """Парсинг содержимого отдельного слайда"""
        # Извлекаем заголовок
        title_match = re.search(r'\*\*Заголовок:\*\*\s*(.+?)(?=\n|$)', content)
        title = title_match.group(1).strip() if title_match else ""
        
        # Извлекаем тело слайда (до секции Иллюстрация)
        body_match = re.search(r'\*\*Тело слайда:\*\*\s*(.+?)(?=\*\*Иллюстрация:\*\*|$)', content, re.DOTALL)
        body = body_match.group(1).strip() if body_match else ""
        
        # Извлекаем описание иллюстрации
        illustration_match = re.search(r'\*\*Иллюстрация:\*\*\s*\[(.+?)\]', content, re.DOTALL)
        illustration_description = illustration_match.group(1).strip() if illustration_match else ""
        
        # Определяем тип слайда
        slide_type = self._determine_slide_type(slide_num, title, content)
        
        return {
            'number': slide_num,
            'title': title,
            'body': body,
            'illustration': illustration_description,
            'type': slide_type
        }

    def _determine_slide_type(self, slide_num, title, content):
        """Определение типа слайда"""
        title_lower = title.lower()
        content_lower = content.lower()
        
        if slide_num == 1 or 'титульный' in title_lower:
            return 'title'
        elif 'перерыв' in title_lower or 'перерыв' in content_lower:
            return 'break'
        elif 'цитата' in title_lower or 'цитата' in content_lower:
            return 'quote'
        else:
            return 'normal'

    def load_template(self):
        """Загрузка шаблона и извлечение изображений"""
        print("Загрузка шаблона...")
        
        try:
            self.template_prs = Presentation(self.template_file)
        except Exception as e:
            print(f"ОШИБКА при загрузке шаблона: {e}")
            sys.exit(1)
            
        if len(self.template_prs.slides) == 0:
            print("ОШИБКА: Шаблон не содержит слайдов")
            sys.exit(1)
            
        # Анализируем первый слайд шаблона
        first_slide = self.template_prs.slides[0]
        
        # Извлекаем все изображения
        for shape in first_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_data = {
                    'image_blob': shape.image.blob,  # Сохраняем бинарные данные
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'rotation': getattr(shape, 'rotation', 0)
                }
                self.template_images.append(image_data)
                
        print(f"Извлечено изображений из шаблона: {len(self.template_images)}")
        print(f"Размер слайда в шаблоне: {self.template_prs.slide_width} x {self.template_prs.slide_height}")
        print(f"Соотношение сторон шаблона: {self.template_prs.slide_width / self.template_prs.slide_height:.2f}")

    def generate_presentation(self):
        """Генерация итоговой презентации с правильным соотношением сторон"""
        print("Генерация презентации...")
        
        # ИСПРАВЛЕНИЕ: Создаем презентацию на основе шаблона для сохранения размеров 16:9
        self.prs = Presentation(self.template_file)
        
        # Проверяем размеры презентации
        print(f"Размер слайда в новой презентации: {self.prs.slide_width} x {self.prs.slide_height}")
        print(f"Соотношение сторон новой презентации: {self.prs.slide_width / self.prs.slide_height:.2f}")
        
        # Удаляем все слайды из шаблона
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
        
        # Генерируем каждый слайд
        for idx, slide_data in enumerate(self.slides_data):
            slide_number = idx + 1  # Номер слайда (начиная с 1)
            self._create_slide(self.prs, slide_data, slide_number)
            
        # Сохраняем презентацию
        try:
            self.prs.save(self.result_file)
            print(f"Презентация сохранена: {self.result_file}")
        except Exception as e:
            print(f"ОШИБКА при сохранении презентации: {e}")
            sys.exit(1)

    def _create_slide(self, prs, slide_data, slide_number):
        """Создание отдельного слайда"""
        # Добавляем пустой слайд
        slide_layout = prs.slide_layouts[6]  # Пустой макет
        slide = prs.slides.add_slide(slide_layout)
        
        # Добавляем все изображения из шаблона
        for img_data in self.template_images:
            # Создаем поток из бинарных данных
            image_stream = BytesIO(img_data['image_blob'])
            
            # Добавляем изображение с точными координатами из шаблона
            pic = slide.shapes.add_picture(
                image_stream,
                img_data['left'],
                img_data['top'],
                img_data['width'],
                img_data['height']
            )
            
            # Применяем поворот если есть
            if img_data['rotation']:
                pic.rotation = img_data['rotation']
        
        # Добавляем текстовое содержимое в зависимости от типа слайда
        if slide_data['type'] in ['title', 'break', 'quote']:
            self._add_special_slide_content(slide, slide_data)
        else:
            self._add_normal_slide_content(slide, slide_data, slide_number)

    def _add_normal_slide_content(self, slide, slide_data, slide_number):
        """Добавление контента для обычного слайда с улучшенным стилем"""
        # Получаем размеры слайда
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        # Определяем позицию текстового блока (чередование)
        is_left_aligned = (slide_number % 2 == 0)  # Четные слайды - слева
        
        # Заголовок (остается на своем месте)
        if slide_data['title']:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(12.5), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.clear()
            title_p = title_frame.paragraphs[0]
            title_p.text = slide_data['title']
            title_p.font.name = 'Montserrat'
            title_p.font.size = Pt(30)
            title_p.font.bold = True  # Жирный шрифт для заголовков
            title_p.alignment = PP_ALIGN.LEFT
            
            # Убираем границы заголовка
            self._remove_shape_borders(title_box)
        
        # Тело слайда с улучшенным стилем
        if slide_data['body']:
            self._add_styled_body_text(slide, slide_data['body'], is_left_aligned)
        
        # Добавляем AI-иллюстрацию если нужно
        if self._should_generate_illustration(slide_number, slide_data):
            image_filename = f"slide_{slide_number:02d}_illustration.png"
            image_path = os.path.join(self.images_dir, image_filename)
            self._add_ai_illustration_to_slide(slide, image_path, slide_data, slide_number)
    
    def _add_styled_body_text(self, slide, body_text, is_left_aligned):
        """Добавляет стилизованный текстовый блок с улучшенным позиционированием"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        # Параметры текстового блока - 50% ширины слайда
        text_width = slide_width * 0.5
        margin = slide_width * 0.05  # 5% отступ от края
        
        # Позиция по горизонтали
        if is_left_aligned:
            left = margin
        else:
            left = slide_width - text_width - margin
        
        # УЛУЧШЕННОЕ вертикальное позиционирование - центр слайда
        title_height = slide_height * 0.15  # Заложим 15% на заголовок
        available_height = slide_height - title_height - (slide_height * 0.1)  # Оставим 10% отступ снизу
        
        # Высота текстового блока - максимум 60% от доступного места
        text_height = min(available_height * 0.6, slide_height * 0.4)
        
        # Центрирование по вертикали
        top = title_height + (available_height - text_height) / 2
        
        # Создаем текстовый блок
        textbox = slide.shapes.add_textbox(left, top, text_width, text_height)
        text_frame = textbox.text_frame
        
        # Настройки фрейма
        text_frame.clear()
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.25)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Текст сверху блока
        
        # Убираем границы
        self._remove_shape_borders(textbox)
        
        # Обработка текста с учетом маркированных списков
        lines = body_text.strip().split('\n')
        first_para = True
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if first_para:
                p = text_frame.paragraphs[0]
                first_para = False
            else:
                p = text_frame.add_paragraph()
            
            # Проверяем, является ли строка элементом списка
            if line.startswith('•'):
                p.text = line[1:].strip()
                p.level = 0  # Уровень списка
            else:
                p.text = line
                p.level = 0
            
            # Форматирование параграфа
            p.font.name = 'Montserrat'
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.LEFT  # Всегда выравнивание по левому краю
            
            # ВАЖНО: Увеличенный межстрочный интервал
            p.line_spacing = 1.5  # Полуторный интервал

    def _add_special_slide_content(self, slide, slide_data):
        """Добавление контента для специальных слайдов с визуальными улучшениями"""
        # Определяем тип слайда для специального форматирования
        slide_type = slide_data.get('type', 'normal')
        
        if slide_type == 'title' and slide_data['number'] == 1:
            # Специальная обработка заглавного слайда
            self._add_title_slide_content(slide, slide_data)
        elif slide_type == 'quote':
            # Специальная обработка цитат
            self._add_quote_slide_content(slide, slide_data)
        else:
            # Обычная обработка для перерывов и прочих спецслайдов
            self._add_generic_special_slide_content(slide, slide_data)
    
    def _add_title_slide_content(self, slide, slide_data):
        """Улучшенное оформление заглавного слайда с переносом строк"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        margin = slide_width * 0.05
        
        # Заголовок - верхняя часть с улучшенным позиционированием
        if slide_data['title']:
            title_width = slide_width * 0.85  # Немного уменьшаем ширину
            title_height = slide_height * 0.35  # Увеличиваем высоту для переноса
            title_top = slide_height * 0.25  # 25% от верха
            title_left = margin + (slide_width * 0.025)  # Центрируем
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.clear()
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            title_frame.word_wrap = True  # Включаем перенос слов
            title_frame.margin_left = Inches(0.2)
            title_frame.margin_right = Inches(0.2)
            title_frame.margin_top = Inches(0.1)
            title_frame.margin_bottom = Inches(0.1)
            
            p = title_frame.paragraphs[0]
            p.text = slide_data['title']
            p.font.name = 'Montserrat'
            p.font.size = Pt(40)  # Уменьшаем размер для лучшего размещения
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.2  # Небольшой межстрочный интервал
            
            self._remove_shape_borders(title_box)
        
        # Подзаголовок - нижняя часть
        if slide_data['body']:
            subtitle_width = slide_width * 0.8  # Меньше ширина
            subtitle_height = slide_height * 0.2
            subtitle_top = slide_height * 0.6  # 60% от верха
            subtitle_left = margin + (slide_width * 0.05)  # Смещение левее
            
            subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            subtitle_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            p = subtitle_frame.paragraphs[0]
            p.text = slide_data['body']
            p.font.name = 'Montserrat'
            p.font.size = Pt(24)  # Меньший размер
            p.font.bold = False
            p.alignment = PP_ALIGN.CENTER
            
            self._remove_shape_borders(subtitle_box)
    
    def _add_quote_slide_content(self, slide, slide_data):
        """Улучшенное оформление слайдов с цитатами"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        margin = slide_width * 0.05
        
        # Обрабатываем текст цитаты
        body_text = slide_data.get('body', '')
        
        # Разделяем цитату и автора (по символу —)
        if '—' in body_text:
            parts = body_text.split('—', 1)
            quote_text = parts[0].strip()
            author_text = '— ' + parts[1].strip()
        else:
            quote_text = body_text
            author_text = ''
        
        text_width = slide_width * 0.9
        text_height = slide_height * 0.6
        left = margin
        top = (slide_height - text_height) / 2
        
        textbox = slide.shapes.add_textbox(left, top, text_width, text_height)
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.5)
        text_frame.margin_right = Inches(0.5)
        text_frame.margin_top = Inches(0.25)
        text_frame.margin_bottom = Inches(0.25)
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Добавляем цитату
        if quote_text:
            p_quote = text_frame.paragraphs[0]
            p_quote.text = quote_text
            p_quote.font.name = 'Montserrat'
            p_quote.font.size = Pt(36)
            p_quote.font.bold = True  # Жирная цитата
            p_quote.alignment = PP_ALIGN.CENTER
        
        # Добавляем автора
        if author_text:
            p_author = text_frame.add_paragraph()
            p_author.text = author_text
            p_author.font.name = 'Montserrat'
            p_author.font.size = Pt(28)
            p_author.font.bold = False
            p_author.font.italic = True  # Курсив для автора
            p_author.alignment = PP_ALIGN.CENTER
            p_author.space_before = Pt(12)  # Отступ перед автором
        
        self._remove_shape_borders(textbox)
    
    def _add_generic_special_slide_content(self, slide, slide_data):
        """Обычное оформление спецслайдов (перерывы)"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        margin = slide_width * 0.05
        
        content_parts = []
        if slide_data['title']:
            content_parts.append(slide_data['title'])
        if slide_data['body']:
            content_parts.append(slide_data['body'])
        
        full_text = "\n\n".join(content_parts)
        
        if full_text:
            text_width = slide_width * 0.9
            text_height = slide_height * 0.6
            left = margin
            top = (slide_height - text_height) / 2
            
            textbox = slide.shapes.add_textbox(left, top, text_width, text_height)
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.5)
            text_frame.margin_right = Inches(0.5)
            text_frame.margin_top = Inches(0.25)
            text_frame.margin_bottom = Inches(0.25)
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = text_frame.paragraphs[0]
            p.text = full_text
            p.font.name = 'Montserrat'
            p.font.size = Pt(40)
            p.font.bold = True  # Жирный шрифт
            p.alignment = PP_ALIGN.CENTER
            
            self._remove_shape_borders(textbox)
    
    def _should_generate_illustration(self, slide_number, slide_data):
        """Определяет, нужно ли генерировать иллюстрацию для слайда"""
        if not self.use_ai_illustrations:
            return False
            
        # Исключаем специальные слайды
        if slide_data['type'] in ['title', 'quote', 'break']:
            return False
            
        # Проверяем интервал
        return slide_number % self.slide_interval == 0
    
    def _generate_image_prompt(self, slide_data):
        """
        Генерирует детальный промпт для DALL-E 3 с использованием лучших практик
        """
        try:
            import requests
            
            url = "https://api.anthropic.com/v1/messages"
            headers = {
                "x-api-key": self.claude_api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json"
            }
            
            # Улучшенный системный промпт с детальными инструкциями
            system_prompt = """Ты эксперт по созданию промптов для DALL-E 3. Создавай промпты, которые генерируют высококачественные, профессиональные изображения для бизнес-презентаций.

КРИТИЧЕСКИ ВАЖНЫЕ ТРЕБОВАНИЯ:
1. ВСЕ изображения должны быть на АБСОЛЮТНО БЕЛОМ ФОНЕ - цвет #FFFFFF в RGB (255, 255, 255)
2. ОБЯЗАТЕЛЬНО укажи "pure white background #FFFFFF" или "absolute white background RGB(255,255,255)"
3. Изображения должны ОСМЫСЛЕННО ДОПОЛНЯТЬ содержание слайда, а не быть абстрактными
4. Используй ЗАГЛАВНЫЕ буквы для важных элементов
5. Добавляй технические детали для фотореализма

ТИПЫ ИЗОБРАЖЕНИЙ в зависимости от контента:
- Для процессов и методологий: создавай СХЕМЫ, ДИАГРАММЫ, FLOWCHARTS
- Для технических тем: визуализируй АРХИТЕКТУРУ, КОМПОНЕНТЫ, СТРУКТУРЫ
- Для концепций: используй МЕТАФОРЫ и ВИЗУАЛЬНЫЕ АНАЛОГИИ
- Для данных: создавай ИНФОГРАФИКУ и ВИЗУАЛИЗАЦИЮ ДАННЫХ

СТИЛИСТИЧЕСКИЕ ТРЕБОВАНИЯ:
- Стиль: современный, минималистичный, профессиональный
- Цветовая схема: яркие акценты на абсолютно белом фоне #FFFFFF
- Качество: ultra high definition, studio lighting
- Техника: isometric view для схем, flat design для иконок, 3D rendering для объектов

СТРУКТУРА ПРОМПТА:
1. Основной объект/концепция
2. Стиль и техника исполнения
3. WHITE BACKGROUND (обязательно)
4. Технические параметры качества
5. Освещение и детализация"""

            # Расширенный пользовательский промпт с контекстом
            illustration_hint = f"\n\nОПИСАНИЕ ИЛЛЮСТРАЦИИ ИЗ ФАЙЛА: {slide_data.get('illustration', 'Не указано')}" if slide_data.get('illustration') else ""
            
            user_prompt = f"""Презентация: Учебные материалы АО "Гознак" по искусственному интеллекту
            
Слайд №{slide_data['number']}
Заголовок: {slide_data['title']}
Содержание: {slide_data['body']}{illustration_hint}

Позиция изображения: {'СПРАВА от текста' if slide_data['number'] % 2 == 0 else 'СЛЕВА от текста'}

Создай промпт для DALL-E 3, который визуализирует ключевую идею этого слайда.

АНАЛИЗ КОНТЕНТА:
1. Определи тип контента (процесс, концепция, данные, методология)
2. Выбери подходящий тип визуализации (схема, диаграмма, метафора, инфографика)
3. Учти позицию изображения при композиции
4. ОБЯЗАТЕЛЬНО используй описание иллюстрации из файла как основу для промпта

ТРЕБОВАНИЯ К ПРОМПТУ:
- Начни с главного объекта/концепции
- Используй ЗАГЛАВНЫЕ буквы для ключевых элементов
- ОБЯЗАТЕЛЬНО укажи "pure white background #FFFFFF" или "absolute white background RGB(255,255,255)"
- Добавь технические параметры: "ultra high definition, studio lighting, professional quality"
- Для схем добавь: "isometric view, 3D rendering, clean design"
- Для концепций: "minimalist style, modern design, bright accent colors"

Промпт должен быть на английском языке, детальным и генерировать изображение, которое РЕАЛЬНО ПОМОГАЕТ понять содержание слайда."""

            data = {
                "model": "claude-3-5-sonnet-20241022",
                "max_tokens": 2048,
                "temperature": 0.9,
                "system": system_prompt,
                "messages": [
                    {
                        "role": "user", 
                        "content": user_prompt
                    }
                ]
            }
            
            response = requests.post(url, headers=headers, json=data, timeout=30)
            
            if response.status_code == 200:
                result = response.json()
                raw_prompt = result['content'][0]['text']
                
                # Постобработка промпта
                enhanced_prompt = self._enhance_dalle_prompt(raw_prompt, slide_data)
                
                # Валидация качества промпта
                validated_prompt = self._validate_prompt_quality(enhanced_prompt)
                
                # Сохраняем промпт в файл
                prompt_filename = f"slide_{slide_data['number']:02d}_prompt.txt"
                prompt_path = os.path.join(self.prompts_dir, prompt_filename)
                
                with open(prompt_path, 'w', encoding='utf-8') as f:
                    f.write(f"Слайд {slide_data['number']}: {slide_data['title']}\n\n")
                    f.write(f"Содержание:\n{slide_data['body']}\n\n")
                    f.write(f"DALL-E Prompt:\n{validated_prompt}")
                
                # Обновляем статистику
                self.execution_stats.increment('prompts_generated')
                self.generation_stats['prompts_generated'] += 1
                
                self.logger.info(f"Промпт для слайда {slide_data['number']} создан: {validated_prompt[:100]}...")
                return validated_prompt
            else:
                self.logger.error(f"Ошибка Claude API: {response.status_code}")
                # Fallback к шаблонам при ошибке API
                return self._generate_fallback_prompt(slide_data)
                
        except Exception as e:
            self.logger.error(f"Ошибка генерации промпта: {e}")
            # Fallback к шаблонам при ошибке
            return self._generate_fallback_prompt(slide_data)
    
    def _generate_fallback_prompt(self, slide_data):
        """
        Генерирует промпт с использованием шаблонов, если Claude API недоступен
        """
        self.logger.info("Использование шаблонов промптов (fallback режим)")
        
        # Используем шаблоны вместо Claude API
        template_prompt = self._select_prompt_template(slide_data)
        
        # Валидация и улучшение
        enhanced_prompt = self._enhance_dalle_prompt(template_prompt, slide_data)
        validated_prompt = self._validate_prompt_quality(enhanced_prompt)
        
        # Сохраняем промпт в файл
        prompt_filename = f"slide_{slide_data['number']:02d}_prompt.txt"
        prompt_path = os.path.join(self.prompts_dir, prompt_filename)
        
        with open(prompt_path, 'w', encoding='utf-8') as f:
            f.write(f"Слайд {slide_data['number']}: {slide_data['title']}\n\n")
            f.write(f"Содержание:\n{slide_data['body']}\n\n")
            f.write(f"TEMPLATE Prompt:\n{validated_prompt}")
        
        # Обновляем статистику
        self.execution_stats.increment('prompts_generated')
        self.generation_stats['prompts_generated'] += 1
        
        self.logger.info(f"Шаблонный промпт для слайда {slide_data['number']} создан: {validated_prompt[:100]}...")
        return validated_prompt
    
    def _enhance_dalle_prompt(self, base_prompt, slide_data):
        """
        Улучшает промпт для DALL-E 3 с учетом лучших практик
        """
        # Добавляем обязательную инструкцию в начало, если её нет
        mandatory_instruction = PromptTemplates.MANDATORY_INSTRUCTION
        if mandatory_instruction not in base_prompt:
            base_prompt = base_prompt.rstrip() + ". " + mandatory_instruction
        
        # Убеждаемся, что есть АБСОЛЮТНО БЕЛЫЙ ФОН (усиленная проверка)
        if "#ffffff" not in base_prompt.lower() and "rgb(255,255,255)" not in base_prompt.lower():
            if "white background" not in base_prompt.lower():
                base_prompt += ", pure white background #FFFFFF RGB(255,255,255)"
            else:
                # Заменяем обычное упоминание на более точное
                base_prompt = base_prompt.replace("white background", "pure white background #FFFFFF")
                base_prompt = base_prompt.replace("WHITE BACKGROUND", "pure white background #FFFFFF RGB(255,255,255)")
        
        # Добавляем технические параметры если их нет
        technical_params = [
            "ultra high definition",
            "studio lighting", 
            "professional quality",
            "sharp focus",
            "highly detailed"
        ]
        
        for param in technical_params:
            if param not in base_prompt.lower():
                base_prompt += f", {param}"
        
        # Специфичные улучшения в зависимости от типа контента
        slide_title_lower = slide_data['title'].lower()
        slide_body_lower = slide_data['body'].lower()
        
        # Для технических/архитектурных тем
        if any(word in slide_title_lower + slide_body_lower for word in 
               ['архитектура', 'структура', 'компонент', 'система', 'модель', 'алгоритм']):
            if "isometric" not in base_prompt.lower():
                base_prompt += ", ISOMETRIC VIEW, 3D diagram, technical illustration"
        
        # Для процессов и методологий
        elif any(word in slide_title_lower + slide_body_lower for word in 
                 ['процесс', 'этап', 'шаг', 'методология', 'подход', 'цикл']):
            if "flowchart" not in base_prompt.lower() and "diagram" not in base_prompt.lower():
                base_prompt += ", FLOWCHART style, connected elements, process visualization"
        
        # Для данных и аналитики
        elif any(word in slide_title_lower + slide_body_lower for word in 
                 ['данные', 'анализ', 'статистика', 'метрика', 'показатель']):
            if "infographic" not in base_prompt.lower():
                base_prompt += ", INFOGRAPHIC style, data visualization, clean charts"
        
        # Для концепций и идей
        elif any(word in slide_title_lower + slide_body_lower for word in 
                 ['идея', 'концепция', 'принцип', 'подход', 'философия']):
            if "metaphor" not in base_prompt.lower() and "concept" not in base_prompt.lower():
                base_prompt += ", CONCEPTUAL illustration, visual metaphor, symbolic representation"
        
        # Добавляем финальную инструкцию для минимизации вмешательства DALL-E
        if len(base_prompt) < 3500:  # Оставляем место для финальной инструкции
            base_prompt += ". I NEED to test how the tool works with extremely simple prompts. DO NOT add any detail, just use it AS-IS"
        
        return base_prompt
    
    def _select_prompt_template(self, slide_data):
        """
        Выбирает подходящий шаблон промпта на основе анализа контента
        """
        title = slide_data['title'].lower()
        body = slide_data['body'].lower()
        content = title + " " + body
        
        # Анализ ключевых слов для выбора шаблона
        if any(word in content for word in ['архитектура', 'структура', 'компонент', 'система']):
            template = PromptTemplates.TECHNICAL_ARCHITECTURE
            concept = self._extract_main_concept(slide_data['title'])
            return template.format(concept=concept)
            
        elif any(word in content for word in ['процесс', 'этап', 'шаг', 'алгоритм', 'последовательность']):
            template = PromptTemplates.PROCESS_FLOW
            process = self._extract_main_concept(slide_data['title'])
            return template.format(process=process)
            
        elif any(word in content for word in ['данные', 'анализ', 'статистика', 'метрика']):
            template = PromptTemplates.DATA_VISUALIZATION
            data_concept = self._extract_main_concept(slide_data['title'])
            return template.format(data_concept=data_concept)
            
        elif any(word in content for word in ['сравнение', 'разница', 'отличие', 'versus']):
            template = PromptTemplates.COMPARISON
            # Извлечение элементов для сравнения
            items = self._extract_comparison_items(slide_data)
            return template.format(item1=items[0], item2=items[1])
            
        elif any(word in content for word in ['история', 'развитие', 'эволюция', 'хронология']):
            template = PromptTemplates.TIMELINE
            topic = self._extract_main_concept(slide_data['title'])
            return template.format(topic=topic)
            
        elif any(word in content for word in ['набор', 'инструменты', 'элементы', 'компоненты']):
            template = PromptTemplates.ICON_SET
            elements = self._extract_main_concept(slide_data['title'])
            return template.format(elements=elements)
            
        else:
            # По умолчанию - концептуальная метафора
            template = PromptTemplates.CONCEPT_METAPHOR
            concept = self._extract_main_concept(slide_data['title'])
            return template.format(concept=concept)
    
    def _extract_main_concept(self, title):
        """
        Извлекает основную концепцию из заголовка слайда
        """
        # Простое извлечение - берем первые значимые слова
        words = title.split()
        # Убираем служебные слова
        stop_words = ['и', 'в', 'на', 'с', 'для', 'по', 'к', 'от', 'до', 'при', 'через', 'между']
        meaningful_words = [word for word in words if word.lower() not in stop_words and len(word) > 2]
        
        # Возвращаем первые 2-3 слова
        if len(meaningful_words) >= 2:
            return ' '.join(meaningful_words[:2])
        elif meaningful_words:
            return meaningful_words[0]
        else:
            return title
    
    def _extract_comparison_items(self, slide_data):
        """
        Извлекает элементы для сравнения из данных слайда
        """
        content = slide_data['title'] + " " + slide_data['body']
        
        # Ищем паттерны сравнения
        if ' vs ' in content.lower():
            parts = content.lower().split(' vs ')
            if len(parts) >= 2:
                return [parts[0].strip(), parts[1].strip()]
        
        if ' против ' in content.lower():
            parts = content.lower().split(' против ')
            if len(parts) >= 2:
                return [parts[0].strip(), parts[1].strip()]
                
        # По умолчанию
        return ["traditional approach", "AI approach"]
    
    def _validate_prompt_quality(self, prompt):
        """
        Проверяет качество промпта перед использованием
        """
        mandatory_instruction = PromptTemplates.MANDATORY_INSTRUCTION
        
        quality_checks = {
            'has_mandatory_instruction': mandatory_instruction in prompt,
            'has_white_background': 'WHITE BACKGROUND' in prompt.upper(),
            'has_4k_quality': '4K' in prompt.upper() or '4k' in prompt.lower(),
            'has_photorealistic': 'photorealistic' in prompt.lower() or 'photorealism' in prompt.lower(),
            'has_quality_params': any(param in prompt.lower() for param in 
                                    ['high definition', 'high quality', 'studio lighting']),
            'has_style_definition': any(style in prompt.lower() for style in 
                                      ['style', 'design', 'illustration', 'visualization']),
            'length_appropriate': 100 < len(prompt) < 3800,
            'has_main_subject': len(prompt.split(',')[0]) > 10,
            'no_text_in_image': 'text:' not in prompt.lower() and 'words:' not in prompt.lower()
        }
        
        failed_checks = [check for check, passed in quality_checks.items() if not passed]
        
        if failed_checks:
            self.logger.warning(f"Промпт не прошел проверки: {failed_checks}")
            
            # Автоматическое исправление критических проблем
            if 'has_mandatory_instruction' in failed_checks:
                prompt = prompt.rstrip() + ". " + mandatory_instruction
            
            if 'has_white_background' in failed_checks:
                prompt += ", WHITE BACKGROUND"
            
            if 'has_quality_params' in failed_checks:
                prompt += ", ultra high definition, studio lighting, professional quality"
        
        return prompt
    
    def _extract_clean_prompt_for_dalle(self, full_prompt):
        """
        Извлекает только английскую часть промпта для отправки в DALL-E API
        Удаляет все метаданные и оставляет только инструкцию на английском
        """
        # Если промпт содержит строки на русском языке (заголовки слайдов),
        # ищем начало английской части
        lines = full_prompt.split('\n')
        english_part = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Пропускаем строки с русскими символами (метаданные)
            if any(ord(char) > 127 for char in line):
                continue
            
            # Пропускаем строки-заголовки типа "Слайд X:", "Содержание:", "DALL-E Prompt:"
            if any(keyword in line for keyword in ['Слайд', 'Содержание:', 'DALL-E Prompt:', 'TEMPLATE Prompt:']):
                continue
                
            # Добавляем только английские строки
            english_part.append(line)
        
        # Объединяем английские части
        clean_prompt = ' '.join(english_part).strip()
        
        # Если не нашли отдельные части, пробуем найти промпт в кавычках
        if not clean_prompt or len(clean_prompt) < 50:
            # Ищем текст в кавычках после "PROMPT:" или "Prompt:"
            import re
            quote_pattern = r'(?:PROMPT|Prompt):\s*["\']([^"\']+)["\']'
            match = re.search(quote_pattern, full_prompt, re.DOTALL | re.IGNORECASE)
            if match:
                clean_prompt = match.group(1).strip()
            else:
                # Fallback: берем весь промпт как есть, но очищаем от явных метаданных
                clean_prompt = full_prompt
                
                # Удаляем строки с русскими символами
                lines = clean_prompt.split('\n')
                english_lines = []
                for line in lines:
                    if not any(ord(char) > 127 for char in line) and line.strip():
                        english_lines.append(line.strip())
                
                if english_lines:
                    clean_prompt = ' '.join(english_lines)
        
        return clean_prompt.strip()
    
    def _generate_image_with_dalle(self, prompt, slide_number):
        """Генерирует изображение с помощью выбранной модели (DALL-E 3 или GPT-Image-1)"""
        # Используем цветной вывод из ColorfulUI.print_image_generation вместо этого
        
        # Извлекаем только английскую часть промпта для OpenAI API
        clean_prompt = self._extract_clean_prompt_for_dalle(prompt)
        
        # Логируем отправляемый промпт
        if self.logger:
            self.logger.info(f"Отправка в {self.image_model} для слайда {slide_number}: {clean_prompt[:100]}...")
        
        # Выбираем метод генерации в зависимости от модели
        if self.image_model == 'gpt-image-1':
            return self._generate_with_gpt_image_1(clean_prompt, slide_number)
        elif self.image_model == 'gemini-2.0-flash':
            return self._generate_with_gemini_flash(clean_prompt, slide_number)
        elif self.image_model == 'imagen-3':
            return self._generate_with_imagen_3(clean_prompt, slide_number)
        else:
            return self._generate_with_dalle_3(clean_prompt, slide_number)
    
    def _generate_with_dalle_3(self, clean_prompt, slide_number):
        """Генерирует изображение с помощью DALL-E 3"""
        url = "https://api.openai.com/v1/images/generations"
        headers = {
            'Authorization': f'Bearer {self.openai_api_key}',
            'Content-Type': 'application/json'
        }
        
        data = {
            'model': 'dall-e-3',
            'prompt': clean_prompt,
            'n': 1,
            'size': '1792x1024',  # 16:9 aspect ratio
            'quality': 'standard',
            'response_format': 'b64_json'
        }
        
        try:
            response = requests.post(url, headers=headers, json=data, timeout=120)
            
            if response.status_code == 200:
                result = response.json()
                image_b64 = result['data'][0]['b64_json']
                
                # Сохраняем изображение
                image_filename = f"slide_{slide_number:02d}_illustration.png"
                image_path = os.path.join(self.images_dir, image_filename)
                
                with open(image_path, 'wb') as f:
                    f.write(base64.b64decode(image_b64))
                
                print(f"✓ Изображение сохранено: {image_filename}")
                
                # Обновляем статистику
                self.execution_stats.increment('images_generated')
                self.generation_stats['images_generated'] += 1
                
                if self.logger:
                    self.logger.info(f"Изображение для слайда {slide_number} успешно создано")
                
                return image_path
                
            else:
                # Детальная информация об ошибке
                error_info = {
                    'slide_number': slide_number,
                    'status_code': response.status_code,
                    'response_text': response.text,
                    'headers': dict(response.headers),
                    'request_url': url,
                    'prompt_preview': prompt[:100] + '...' if len(prompt) > 100 else prompt
                }
                
                self._log_error("DALL-E API Error", error_info)
                
                # Обновляем статистику
                self.execution_stats.increment('images_failed')
                self.execution_stats.increment('total_errors')
                self.generation_stats['images_failed'] += 1
                
                print(f"\n❌ Ошибка DALL-E API для слайда {slide_number}:")
                print(f"   Статус: {response.status_code}")
                print(f"   Endpoint: {url}")
                
                if response.status_code == 401:
                    print("   → Ошибка авторизации: проверьте OpenAI API ключ")
                elif response.status_code == 429:
                    print("   → Превышен лимит запросов или недостаточно кредитов")
                elif response.status_code == 400:
                    print("   → Ошибка в промпте или параметрах")
                    print(f"   → Ответ: {response.text[:200]}...")
                else:
                    print(f"   → HTTP {response.status_code}: {response.text[:200]}...")
                    
                return None
                
        except requests.exceptions.Timeout:
            error_msg = f"Таймаут при генерации изображения для слайда {slide_number}"
            print(f"❌ {error_msg}")
            self._log_error("DALL-E API Timeout", {'slide_number': slide_number, 'error': error_msg})
            
            # Обновляем статистику
            self.execution_stats.increment('images_failed')
            self.execution_stats.increment('total_errors')
            self.generation_stats['images_failed'] += 1
            return None
            
        except requests.exceptions.ConnectionError:
            error_msg = f"Ошибка соединения с DALL-E API для слайда {slide_number}"
            print(f"❌ {error_msg}")
            self._log_error("DALL-E API Connection Error", {'slide_number': slide_number, 'error': error_msg})
            
            # Обновляем статистику
            self.execution_stats.increment('images_failed')
            self.execution_stats.increment('total_errors')
            self.generation_stats['images_failed'] += 1
            return None
            
        except Exception as e:
            error_msg = f"Неожиданная ошибка при генерации изображения для слайда {slide_number}: {type(e).__name__}: {str(e)}"
            print(f"❌ {error_msg}")
            self._log_error("Unexpected DALL-E API Error", {
                'slide_number': slide_number, 
                'error_type': type(e).__name__, 
                'error_message': str(e)
            })
            
            # Обновляем статистику
            self.execution_stats.increment('images_failed')
            self.execution_stats.increment('total_errors')
            self.generation_stats['images_failed'] += 1
            return None
    
    def _generate_with_gpt_image_1(self, clean_prompt, slide_number):
        """Генерирует изображение с помощью GPT-Image-1"""
        try:
            from openai import OpenAI
            
            # Инициализируем клиент OpenAI
            client = OpenAI(api_key=self.openai_api_key)
            
            # Параметры для GPT-Image-1 (оптимизированы для презентаций)
            generation_params = {
                'model': 'gpt-image-1',
                'prompt': clean_prompt,
                'n': 1,
                'size': '1536x1024',  # Landscape format 3:2 (близко к 16:9)
                'quality': 'high',  # Высокое качество
                'background': 'opaque',  # Непрозрачный фон (белый)
                'output_format': 'png',  # PNG формат для лучшего качества
                'output_compression': 90,  # Умеренное сжатие для баланса качества/размера
                'moderation': 'auto'  # Автоматическая модерация
            }
            
            if self.logger:
                self.logger.info(f"GPT-Image-1 запрос для слайда {slide_number}: {generation_params}")
            
            # Генерируем изображение
            response = client.images.generate(**generation_params)
            
            # Получаем base64 данные
            image_b64 = response.data[0].b64_json
            
            # Сохраняем изображение
            image_filename = f"slide_{slide_number:02d}_illustration.png"
            image_path = os.path.join(self.images_dir, image_filename)
            
            # Декодируем и сохраняем
            image_bytes = base64.b64decode(image_b64)
            with open(image_path, 'wb') as f:
                f.write(image_bytes)
            
            print(f"✓ Изображение GPT-Image-1 сохранено: {image_filename}")
            
            # Логируем использование токенов если доступно
            if hasattr(response, 'usage') and response.usage:
                if self.logger:
                    self.logger.info(f"GPT-Image-1 токены для слайда {slide_number}: {response.usage}")
                print(f"   Использовано токенов: {response.usage.total_tokens}")
            
            # Обновляем статистику
            self.execution_stats.increment('images_generated')
            self.generation_stats['images_generated'] += 1
            
            if self.logger:
                self.logger.info(f"Изображение GPT-Image-1 для слайда {slide_number} успешно создано")
            
            return image_path
            
        except ImportError:
            error_msg = "Для использования GPT-Image-1 необходимо установить библиотеку openai: pip install openai>=1.0.0"
            print(f"❌ {error_msg}")
            self._log_error("OpenAI Library Error", {'slide_number': slide_number, 'error': error_msg})
            
            # Обновляем статистику
            self.execution_stats.increment('images_failed')
            self.execution_stats.increment('total_errors')
            self.generation_stats['images_failed'] += 1
            return None
            
        except Exception as e:
            error_msg = f"Ошибка GPT-Image-1 для слайда {slide_number}: {type(e).__name__}: {str(e)}"
            print(f"❌ {error_msg}")
            self._log_error("GPT-Image-1 API Error", {
                'slide_number': slide_number, 
                'error_type': type(e).__name__, 
                'error_message': str(e)
            })
            
            # Обновляем статистику
            self.execution_stats.increment('images_failed')
            self.execution_stats.increment('total_errors')
            self.generation_stats['images_failed'] += 1
            return None
    
    def _generate_with_imagen_3(self, clean_prompt, slide_number):
        """Генерирует изображение с помощью Google Imagen 3"""
        try:
            from google import genai
            from google.genai import types
            from PIL import Image
            from io import BytesIO
            
            # Проверяем наличие API ключа
            if not hasattr(self, 'gemini_api_key') or not self.gemini_api_key:
                error_msg = "Для использования Imagen 3 необходимо настроить GEMINI_API_KEY"
                print(f"❌ {error_msg}")
                self._log_error("Imagen 3 API Key Error", {'slide_number': slide_number, 'error': error_msg})
                
                # Обновляем статистику
                self.execution_stats.increment('images_failed')
                self.execution_stats.increment('total_errors')
                self.generation_stats['images_failed'] += 1
                return None
            
            # Инициализируем клиент Gemini
            client = genai.Client(api_key=self.gemini_api_key)
            
            # Параметры для Imagen 3 (оптимизированы для презентаций)
            config = types.GenerateImagesConfig(
                number_of_images=1,
                aspect_ratio="16:9",  # Оптимально для презентаций
                person_generation="allow_adult"  # Разрешаем взрослых людей
            )
            
            if self.logger:
                self.logger.info(f"Imagen 3 запрос для слайда {slide_number}: {clean_prompt[:100]}...")
            
            # Генерируем изображение через Imagen 3
            response = client.models.generate_images(
                model='imagen-3.0-generate-002',
                prompt=clean_prompt,
                config=config
            )
            
            if not response.generated_images:
                raise Exception("Не получено изображений от Imagen 3")
            
            # Получаем первое изображение
            generated_image = response.generated_images[0]
            image_bytes = generated_image.image.image_bytes
            
            # Создаем директорию если не существует
            if not os.path.exists(self.images_dir):
                os.makedirs(self.images_dir)
            
            # Сохраняем изображение
            image_filename = f"slide_{slide_number:02d}_illustration.png"
            image_path = os.path.join(self.images_dir, image_filename)
            
            # Конвертируем байты в изображение и сохраняем
            image = Image.open(BytesIO(image_bytes))
            image.save(image_path, 'PNG')
            
            print(f"✓ Изображение Imagen 3 сохранено: {image_filename}")
            
            # Обновляем статистику
            self.execution_stats.increment('images_generated')
            self.generation_stats['images_generated'] += 1
            
            if self.logger:
                self.logger.info(f"Изображение Imagen 3 для слайда {slide_number} успешно создано")
            
            return image_path
            
        except ImportError:
            error_msg = "Для использования Imagen 3 необходимо установить библиотеку google-genai: pip install google-genai"
            print(f"❌ {error_msg}")
            self._log_error("Imagen 3 Library Error", {'slide_number': slide_number, 'error': error_msg})
            
            # Обновляем статистику
            self.execution_stats.increment('images_failed')
            self.execution_stats.increment('total_errors')
            self.generation_stats['images_failed'] += 1
            return None
            
        except Exception as e:
            error_msg = f"Ошибка генерации с Imagen 3: {str(e)}"
            print(f"❌ {error_msg}")
            self._log_error("Imagen 3 Error", {
                'slide_number': slide_number, 
                'error_type': type(e).__name__, 
                'error_message': str(e)
            })
            
            # Обновляем статистику
            self.execution_stats.increment('images_failed')
            self.execution_stats.increment('total_errors')
            self.generation_stats['images_failed'] += 1
            return None
    
    def _generate_with_gemini_flash(self, clean_prompt, slide_number):
        """
        DEPRECATED: Gemini 2.0 Flash image generation is not available.
        Redirects to Imagen 3 for image generation.
        """
        print("⚠️  Gemini 2.0 Flash image generation is not available in the current API")
        print("🔄 Redirecting to Imagen 3 for image generation...")
        
        if self.logger:
            self.logger.warning(f"Gemini 2.0 Flash не доступен для слайда {slide_number}, используем Imagen 3")
        
        # Redirect to Imagen 3 which actually works
        return self._generate_with_imagen_3(clean_prompt, slide_number)
    
    def _add_ai_illustration_to_slide(self, slide, image_path, slide_data, slide_number):
        """Добавляет AI-иллюстрацию на слайд с адаптивным позиционированием"""
        if not os.path.exists(image_path):
            return
            
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        # Определяем позицию изображения (противоположную тексту)
        is_left_aligned = (slide_number % 2 == 0)  # Четные слайды - текст слева
        
        # Размеры изображения - 40% ширины слайда
        img_width = slide_width * 0.4
        img_height = img_width * 0.5625  # Соотношение 16:9
        
        # Позиционирование
        margin = slide_width * 0.05
        
        if is_left_aligned:
            # Текст слева, изображение справа
            img_left = slide_width - img_width - margin
        else:
            # Текст справа, изображение слева
            img_left = margin
        
        # Вертикальное центрирование с учетом заголовка
        title_height = slide_height * 0.15
        available_height = slide_height - title_height - (slide_height * 0.1)
        img_top = title_height + (available_height - img_height) / 2
        
        try:
            # Добавляем изображение
            pic = slide.shapes.add_picture(
                image_path,
                img_left,
                img_top,
                img_width,
                img_height
            )
            
            # Убираем границы
            self._remove_shape_borders(pic)
            
            # Обновляем статистику успешной вставки
            self.execution_stats.increment('images_inserted')
            
        except Exception as e:
            print(f"Ошибка при добавлении изображения: {e}")
            self.execution_stats.increment('total_errors')
    
    def _process_ai_illustrations(self):
        """СТРОГАЯ генерация AI-иллюстраций с обязательными проверками"""
        if not self.use_ai_illustrations:
            return True
            
        print(f"\n=== СТРОГАЯ ГЕНЕРАЦИЯ AI-ИЛЛЮСТРАЦИЙ ===")
        
        # Контрольная точка 1: Проверка валидации API
        if not self.checkpoints.validate_checkpoint('api_validation'):
            print("⛔ ОСТАНОВКА: API ключи не прошли валидацию")
            self._cleanup_partial_results()
            sys.exit(1)
        
        # Определяем слайды для обработки
        slides_to_process = []
        for slide_data in self.slides_data:
            if self._should_generate_illustration(slide_data['number'], slide_data):
                slides_to_process.append(slide_data)
        
        if not slides_to_process:
            print("Нет слайдов для генерации иллюстраций")
            return True
            
        # Инициализируем статистику
        total_slides = len(slides_to_process)
        self.execution_stats.set('slides_to_process', total_slides)
        self.generation_stats['total_slides'] = total_slides
        
        print(f"📊 Запланировано к обработке: {total_slides} слайдов")
        print(f"📋 Слайды: {[s['number'] for s in slides_to_process]}")
        
        # ЭТАП 1: Генерация промптов
        ColorfulUI.print_step(1, "Генерация промптов с Claude API", emoji='magic')
        prompts_success = self._generate_all_prompts(slides_to_process)
        
        if not prompts_success:
            print("\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Генерация промптов провалилась")
            self._cleanup_partial_results()
            sys.exit(2)
        
        # Контрольная точка 2: Проверка промптов
        if not self.checkpoints.validate_checkpoint('prompts_generation'):
            print("\n⛔ ОСТАНОВКА: Недостаточно успешных промптов")
            self._cleanup_partial_results()
            sys.exit(3)
        
        # ЭТАП 2: Генерация изображений
        print(f"\n🔸 ЭТАП 2: Генерация изображений с DALL-E 3")
        images_success = self._generate_all_images()
        
        if not images_success:
            print("\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Генерация изображений провалилась")
            self._cleanup_partial_results()
            sys.exit(4)
        
        # Контрольная точка 3: Проверка изображений
        if not self.checkpoints.validate_checkpoint('images_generation'):
            print("\n⛔ ОСТАНОВКА: Недостаточно успешных изображений")
            self._cleanup_partial_results()
            sys.exit(5)
        
        # ЭТАП 3: Вставка изображений в презентацию
        print(f"\n🔸 ЭТАП 3: Вставка изображений в презентацию")
        update_success = self._update_presentation_with_images(slides_to_process)
        
        if not update_success:
            print("\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Не удалось вставить изображения в презентацию")
            self._cleanup_partial_results()
            sys.exit(6)
        
        # Контрольная точка 4: Проверка обновления презентации
        if not self.checkpoints.validate_checkpoint('presentation_update'):
            print("\n⛔ ОСТАНОВКА: Ошибка валидации обновления презентации")
            self._cleanup_partial_results()
            sys.exit(6)
        
        print(f"\n✅ Все этапы генерации AI-иллюстраций завершены успешно")
        return True
    
    def _process_ai_illustrations_parallel(self):
        """ПАРАЛЛЕЛЬНАЯ генерация AI-иллюстраций с оптимизацией производительности"""
        if not self.use_ai_illustrations:
            return True
            
        print(f"\n=== ПАРАЛЛЕЛЬНАЯ ГЕНЕРАЦИЯ AI-ИЛЛЮСТРАЦИЙ ===")
        
        # Контрольная точка 1: Проверка валидации API
        if not self.checkpoints.validate_checkpoint('api_validation'):
            print("⛔ ОСТАНОВКА: API ключи не прошли валидацию")
            self._cleanup_partial_results()
            sys.exit(1)
        
        # Определяем слайды для обработки
        slides_to_process = []
        for slide_data in self.slides_data:
            if self._should_generate_illustration(slide_data['number'], slide_data):
                slides_to_process.append(slide_data)
        
        if not slides_to_process:
            print("Нет слайдов для генерации иллюстраций")
            return True
            
        # Инициализируем статистику
        total_slides = len(slides_to_process)
        self.execution_stats.set('slides_to_process', total_slides)
        self.generation_stats['total_slides'] = total_slides
        
        print(f"📊 Запланировано к обработке: {total_slides} слайдов")
        print(f"📋 Слайды: {[s['number'] for s in slides_to_process]}")
        
        # Инициализируем очереди для параллельной обработки
        prompt_queue = queue.Queue()
        image_queue = queue.Queue()
        results = {}
        
        # Запускаем параллельную обработку
        success = self._run_parallel_generation(slides_to_process, prompt_queue, image_queue, results)
        
        if not success:
            print("\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Параллельная генерация провалилась")
            self._cleanup_partial_results()
            sys.exit(2)
        
        # ЭТАП 3: Вставка изображений в презентацию
        print(f"\n🔸 ЭТАП 3: Вставка изображений в презентацию")
        update_success = self._update_presentation_with_images(slides_to_process)
        
        if not update_success:
            print("\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Не удалось вставить изображения в презентацию")
            self._cleanup_partial_results()
            sys.exit(6)
        
        # Контрольная точка 4: Проверка обновления презентации
        if not self.checkpoints.validate_checkpoint('presentation_update'):
            print("\n⛔ ОСТАНОВКА: Ошибка валидации обновления презентации")
            self._cleanup_partial_results()
            sys.exit(6)
        
        print(f"\n✅ Все этапы параллельной генерации AI-иллюстраций завершены успешно")
        return True
    
    def _run_parallel_generation(self, slides_to_process, prompt_queue, image_queue, results):
        """Запускает параллельную генерацию промптов и изображений"""
        
        # Блокировки для потокобезопасности
        prompt_lock = threading.Lock()
        image_lock = threading.Lock()
        results_lock = threading.Lock()
        
        # Счетчики для отслеживания прогресса
        prompt_progress = {'completed': 0, 'total': len(slides_to_process)}
        image_progress = {'completed': 0, 'total': 0}
        
        def prompt_worker():
            """Воркер для генерации промптов"""
            while True:
                try:
                    slide_data = prompt_queue.get(timeout=1)
                    if slide_data is None:  # Сигнал завершения
                        break
                    
                    # Генерируем промпт
                    with prompt_lock:
                        self.execution_stats.increment('prompts_attempted')
                    
                    dalle_prompt = self._generate_image_prompt(slide_data)
                    
                    if dalle_prompt:
                        with prompt_lock:
                            self.execution_stats.increment('prompts_generated')
                            prompt_progress['completed'] += 1
                        
                        # Сразу отправляем задачу на генерацию изображения
                        image_queue.put({
                            'slide_data': slide_data,
                            'prompt': dalle_prompt
                        })
                        
                        with image_lock:
                            image_progress['total'] += 1
                            
                        print(f"✓ Промпт готов для слайда {slide_data['number']}, отправлен на генерацию изображения")
                    else:
                        with prompt_lock:
                            self.execution_stats.increment('prompts_failed')
                            print(f"❌ Ошибка генерации промпта для слайда {slide_data['number']}")
                    
                    prompt_queue.task_done()
                    time.sleep(1)  # Задержка между запросами к Claude API
                    
                except queue.Empty:
                    continue
                except Exception as e:
                    print(f"❌ Ошибка в prompt_worker: {e}")
                    break
        
        def image_worker():
            """Воркер для генерации изображений"""
            while True:
                try:
                    prompt_data = image_queue.get(timeout=5)
                    if prompt_data is None:  # Сигнал завершения
                        break
                    
                    slide_data = prompt_data['slide_data']
                    prompt = prompt_data['prompt']
                    slide_number = slide_data['number']
                    
                    with image_lock:
                        self.execution_stats.increment('images_attempted')
                    
                    # Генерируем изображение
                    image_path = self._generate_image_with_dalle(prompt, slide_number)
                    
                    if image_path:
                        with image_lock:
                            self.execution_stats.increment('images_generated')
                            image_progress['completed'] += 1
                        
                        with results_lock:
                            results[slide_number] = {
                                'slide_data': slide_data,
                                'prompt': prompt,
                                'image_path': image_path
                            }
                        
                        print(f"✓ Изображение готово для слайда {slide_number}")
                    else:
                        with image_lock:
                            self.execution_stats.increment('images_failed')
                            print(f"❌ Ошибка генерации изображения для слайда {slide_number}")
                    
                    image_queue.task_done()
                    time.sleep(2)  # Задержка между запросами к DALL-E API
                    
                except queue.Empty:
                    continue
                except Exception as e:
                    print(f"❌ Ошибка в image_worker: {e}")
                    break
        
        # Заполняем очередь промптов
        for slide_data in slides_to_process:
            prompt_queue.put(slide_data)
        
        # Запускаем воркеры
        print(f"\n🔸 ПАРАЛЛЕЛЬНАЯ ГЕНЕРАЦИЯ: Запуск воркеров")
        
        # 1 воркер для промптов (ограничение Claude API)
        prompt_thread = threading.Thread(target=prompt_worker, name="PromptWorker")
        prompt_thread.daemon = True
        prompt_thread.start()
        
        # 2 воркера для изображений (можно больше запросов к DALL-E)
        image_threads = []
        for i in range(2):
            thread = threading.Thread(target=image_worker, name=f"ImageWorker-{i+1}")
            thread.daemon = True
            thread.start()
            image_threads.append(thread)
        
        # Ждем завершения генерации промптов
        print(f"📝 Генерация промптов...")
        while prompt_progress['completed'] < prompt_progress['total']:
            time.sleep(1)
            print(f"   Промпты: {prompt_progress['completed']}/{prompt_progress['total']}")
        
        # Проверяем успешность промптов
        prompts_success_rate = self.execution_stats.get('prompts_generated') / len(slides_to_process)
        if prompts_success_rate < 0.8:
            print(f"❌ Успешность промптов {prompts_success_rate:.1%} ниже требуемых 80%")
            return False
        
        print(f"✅ Генерация промптов завершена: {prompts_success_rate:.1%} успешности")
        
        # Ждем завершения генерации изображений
        print(f"🎨 Генерация изображений...")
        while image_progress['completed'] < image_progress['total']:
            time.sleep(2)
            print(f"   Изображения: {image_progress['completed']}/{image_progress['total']}")
        
        # Останавливаем воркеры
        prompt_queue.put(None)
        for _ in image_threads:
            image_queue.put(None)
        
        # Ждем завершения потоков
        prompt_thread.join(timeout=10)
        for thread in image_threads:
            thread.join(timeout=10)
        
        # Проверяем успешность изображений
        images_success_rate = self.execution_stats.get('images_generated') / image_progress['total']
        if images_success_rate < 0.8:
            print(f"❌ Успешность изображений {images_success_rate:.1%} ниже требуемых 80%")
            return False
        
        print(f"✅ Генерация изображений завершена: {images_success_rate:.1%} успешности")
        
        # Сохраняем результаты для следующего этапа
        self.parallel_results = results
        
        return True
    
    def _generate_all_prompts(self, slides_to_process):
        """Генерирует промпты для всех слайдов с контролем качества"""
        total_slides = len(slides_to_process)
        
        # Создаем цветной прогресс-бар
        progress_bar = ColorfulUI.create_progress_bar(total_slides, "🧠 Промпты:", "bright_cyan")
        successful_prompts = []
        
        for i, slide_data in enumerate(slides_to_process):
            self.execution_stats.increment('prompts_attempted')
            
            # Цветное сообщение о генерации промпта
            ColorfulUI.print_prompt_generation(slide_data['number'], total_slides)
            
            # Генерируем промпт
            dalle_prompt = self._generate_image_prompt(slide_data)
            if dalle_prompt:
                successful_prompts.append({
                    'slide_data': slide_data,
                    'prompt': dalle_prompt
                })
            
            # Обновляем прогресс-бар
            progress_bar.update(i + 1)
            
            # Задержка между запросами
            time.sleep(1)
        
        # Проверяем результат
        success_rate = self.execution_stats.get('prompts_generated') / len(slides_to_process)
        print(f"\n📊 Результат генерации промптов:")
        print(f"   Успешно: {self.execution_stats.get('prompts_generated')}")
        print(f"   Неудачно: {self.execution_stats.get('prompts_failed')}")
        print(f"   Успешность: {success_rate:.1%}")
        
        if success_rate < 0.8:
            print(f"❌ Успешность {success_rate:.1%} ниже требуемых 80%")
            return False
        
        # Сохраняем успешные промпты для следующего этапа
        self.successful_prompts = successful_prompts
        return True
    
    def _generate_all_images(self):
        """Генерирует изображения для всех промптов с контролем качества"""
        if not hasattr(self, 'successful_prompts'):
            print("❌ Нет успешных промптов для генерации изображений")
            return False
        
        prompts_count = len(self.successful_prompts)
        
        # Создаем цветной прогресс-бар для изображений
        progress_bar = ColorfulUI.create_progress_bar(prompts_count, "🪄 Изображения:", "bright_magenta")
        
        for i, prompt_data in enumerate(self.successful_prompts):
            slide_data = prompt_data['slide_data']
            prompt = prompt_data['prompt']
            slide_number = slide_data['number']
            
            self.execution_stats.increment('images_attempted')
            
            # Цветное сообщение о генерации изображения
            ColorfulUI.print_image_generation(slide_number, prompts_count, self.image_model)
            
            # Генерируем изображение
            image_path = self._generate_image_with_dalle(prompt, slide_number)
            
            # Обновляем прогресс-бар
            progress_bar.update(i + 1)
            
            # Задержка между запросами
            time.sleep(2)
        
        # Проверяем результат
        success_rate = self.execution_stats.get('images_generated') / prompts_count
        print(f"\n📊 Результат генерации изображений:")
        print(f"   Успешно: {self.execution_stats.get('images_generated')}")
        print(f"   Неудачно: {self.execution_stats.get('images_failed')}")
        print(f"   Успешность: {success_rate:.1%}")
        
        if success_rate < 0.8:
            print(f"❌ Успешность {success_rate:.1%} ниже требуемых 80%")
            return False
        
        return True
    
    def _cleanup_partial_results(self, save_history=True):
        """
        Очищает частично созданные файлы
        
        Args:
            save_history: Если True, сначала сохраняет в историю
        """
        print("\n🧹 Обработка результатов генерации...")
        
        # Сначала пытаемся сохранить в историю
        if save_history and (os.path.exists(self.prompts_dir) or os.path.exists(self.images_dir)):
            history_saved = self._save_generation_history()
            
            if history_saved:
                print("📁 Файлы сохранены в историю перед очисткой")
        
        # Теперь очищаем временные директории
        try:
            import shutil
            
            # Удаляем только временные файлы, НЕ трогая историю
            if os.path.exists(self.prompts_dir):
                shutil.rmtree(self.prompts_dir)
                print("✓ Временные промпты удалены")
                
            if os.path.exists(self.images_dir):
                shutil.rmtree(self.images_dir)
                print("✓ Временные изображения удалены")
                
        except Exception as e:
            print(f"⚠️  Ошибка при очистке: {e}")

    def _update_presentation_with_images(self, slides_to_process):
        """
        Обновляет презентацию, вставляя AI-изображения
        КРИТИЧЕСКИ ВАЖНО: этот метод должен действительно вставлять изображения
        """
        try:
            print("📎 Вставка изображений в презентацию...")
            
            # Проверяем, что презентация уже создана
            if not os.path.exists(self.result_file):
                print(f"❌ Презентация не найдена: {self.result_file}")
                return False
            
            # Загружаем существующую презентацию
            prs = Presentation(self.result_file)
            images_inserted = 0
            
            # Подготавливаем данные о сгенерированных изображениях
            generated_images = {}
            for slide_data in slides_to_process:
                slide_num = slide_data['number']
                image_filename = f"slide_{slide_num:02d}_illustration.png"
                image_path = os.path.join(self.images_dir, image_filename)
                
                if os.path.exists(image_path):
                    generated_images[slide_num] = {
                        'path': image_path,
                        'filename': image_filename,
                        'slide_data': slide_data
                    }
            
            print(f"Найдено {len(generated_images)} изображений для вставки")
            
            # Для каждого сгенерированного изображения
            for slide_num, image_info in generated_images.items():
                try:
                    # Получаем слайд (нумерация с 0)
                    slide_index = slide_num - 1
                    if slide_index >= len(prs.slides):
                        print(f"❌ Слайд {slide_num} не существует в презентации")
                        continue
                    
                    slide = prs.slides[slide_index]
                    image_path = image_info['path']
                    
                    # Определяем позицию для изображения
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    
                    # Размеры изображения (40% ширины слайда)
                    image_width = int(slide_width * 0.4)
                    
                    # Позиционирование (противоположно тексту)
                    margin = int(slide_width * 0.05)
                    if slide_num % 2 == 0:  # Четный - текст слева, изображение справа
                        image_left = slide_width - image_width - margin
                    else:  # Нечетный - текст справа, изображение слева
                        image_left = margin
                    
                    # Вертикальное центрирование с учетом заголовка
                    title_height = int(slide_height * 0.15)
                    available_height = slide_height - title_height - int(slide_height * 0.1)
                    
                    # Рассчитываем высоту изображения (16:9 aspect ratio)
                    image_height = int(image_width * 0.5625)
                    
                    # Если изображение слишком высокое, уменьшаем
                    max_height = int(available_height * 0.8)
                    if image_height > max_height:
                        image_height = max_height
                        image_width = int(image_height / 0.5625)
                    
                    # Центрируем по вертикали
                    image_top = title_height + (available_height - image_height) // 2
                    
                    # КРИТИЧЕСКИ ВАЖНО: Реально добавляем изображение
                    picture = slide.shapes.add_picture(
                        image_path,
                        left=image_left,
                        top=image_top,
                        width=image_width,
                        height=image_height
                    )
                    
                    images_inserted += 1
                    print(f"✓ Изображение вставлено на слайд {slide_num}")
                    
                    if self.logger:
                        self.logger.info(f"Изображение вставлено на слайд {slide_num}: {image_path}")
                    
                except Exception as e:
                    print(f"❌ Ошибка вставки изображения на слайд {slide_num}: {e}")
                    if self.logger:
                        self.logger.error(f"Ошибка вставки изображения на слайд {slide_num}: {e}")
                        import traceback
                        self.logger.error(traceback.format_exc())
                    continue
            
            # Сохраняем обновленную презентацию
            if images_inserted > 0:
                # Новое имя файла с суффиксом
                illustrated_path = os.path.join(self.result_dir, "RWTech_Universal_Presentation_Illustrated.pptx")
                prs.save(illustrated_path)
                print(f"✓ Презентация с иллюстрациями сохранена: {illustrated_path}")
                
                # Обновляем статистику
                self.execution_stats.set('images_inserted', images_inserted)
                
                if self.logger:
                    self.logger.info(f"Презентация с иллюстрациями сохранена: {illustrated_path}")
                
                return True
            else:
                print("❌ Ни одно изображение не было вставлено")
                return False
                
        except Exception as e:
            print(f"❌ Критическая ошибка при обновлении презентации: {e}")
            if self.logger:
                self.logger.error(f"Критическая ошибка при обновлении презентации: {e}")
                import traceback
                self.logger.error(traceback.format_exc())
            return False

    def _save_generation_history(self):
        """
        Сохраняет промпты и изображения в историю для возможного повторного использования
        """
        try:
            import shutil
            
            # Создаем timestamp для уникальности
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            # Создаем директорию истории
            history_base = os.path.join(self.base_path, "history")
            os.makedirs(history_base, exist_ok=True)
            
            # Создаем директорию для текущей генерации
            generation_dir = os.path.join(history_base, f"generation_{timestamp}")
            os.makedirs(generation_dir, exist_ok=True)
            
            # Сохраняем метаданные генерации
            metadata = {
                'timestamp': timestamp,
                'date': datetime.now().isoformat(),
                'total_slides': len(self.slides_data) if self.slides_data else 60,
                'ai_mode': self.use_ai_illustrations,
                'slide_interval': self.slide_interval if hasattr(self, 'slide_interval') else 5,
                'prompts_attempted': self.execution_stats.get('prompts_attempted', 0),
                'prompts_generated': self.execution_stats.get('prompts_generated', 0),
                'images_attempted': self.execution_stats.get('images_attempted', 0),
                'images_generated': self.execution_stats.get('images_generated', 0),
                'images_inserted': self.execution_stats.get('images_inserted', 0),
                'success': self.execution_stats.get('images_inserted', 0) > 0
            }
            
            metadata_path = os.path.join(generation_dir, 'metadata.json')
            with open(metadata_path, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, ensure_ascii=False, indent=2)
            
            files_saved = False
            
            # Копируем промпты если они есть
            if os.path.exists(self.prompts_dir) and os.listdir(self.prompts_dir):
                prompts_history = os.path.join(generation_dir, "prompts")
                shutil.copytree(self.prompts_dir, prompts_history)
                print(f"📁 Промпты сохранены в историю: {prompts_history}")
                files_saved = True
            
            # Копируем изображения если они есть
            if os.path.exists(self.images_dir) and os.listdir(self.images_dir):
                images_history = os.path.join(generation_dir, "images")
                shutil.copytree(self.images_dir, images_history)
                print(f"📁 Изображения сохранены в историю: {images_history}")
                files_saved = True
            
            # Копируем финальную презентацию если она есть
            illustrated_path = os.path.join(self.result_dir, "RWTech_Universal_Presentation_Illustrated.pptx")
            if os.path.exists(illustrated_path):
                shutil.copy2(
                    illustrated_path, 
                    os.path.join(generation_dir, f"presentation_{timestamp}.pptx")
                )
                print(f"📁 Презентация сохранена в историю")
                files_saved = True
            elif os.path.exists(self.result_file):
                # Сохраняем обычную презентацию если иллюстрированной нет
                shutil.copy2(
                    self.result_file,
                    os.path.join(generation_dir, f"presentation_{timestamp}.pptx")
                )
                print(f"📁 Презентация сохранена в историю")
                files_saved = True
            
            # Копируем лог файл
            if hasattr(self, 'logger') and self.logger and hasattr(self.logger, 'handlers'):
                for handler in self.logger.handlers:
                    if hasattr(handler, 'baseFilename'):
                        log_file = handler.baseFilename
                        if os.path.exists(log_file):
                            shutil.copy2(log_file, os.path.join(generation_dir, "generation.log"))
                            break
            
            if files_saved:
                print(f"\n✅ История генерации сохранена: {generation_dir}")
                if self.logger:
                    self.logger.info(f"История генерации сохранена: {generation_dir}")
            else:
                # Удаляем пустую директорию
                os.rmdir(generation_dir)
                print(f"⚠️  Нет файлов для сохранения в историю")
            
            return files_saved
            
        except Exception as e:
            print(f"⚠️  Ошибка сохранения истории: {e}")
            if self.logger:
                self.logger.error(f"Ошибка сохранения истории: {e}")
            return False

    def validate_final_result(self):
        """Комплексная проверка финального результата"""
        print(f"\n🔍 ФИНАЛЬНАЯ ВАЛИДАЦИЯ РЕЗУЛЬТАТА")
        
        validation_results = {
            'file_exists': False,
            'file_size_ok': False,
            'slide_count_correct': False,
            'images_inserted': False,
            'no_corruption': False
        }
        
        # Определяем какой файл проверять
        illustrated_path = os.path.join(self.result_dir, "RWTech_Universal_Presentation_Illustrated.pptx")
        final_file = illustrated_path if (self.use_ai_illustrations and os.path.exists(illustrated_path)) else self.result_file
        
        # 1. Проверка существования файла
        if not os.path.exists(final_file):
            print(f"❌ Финальный файл презентации не найден: {final_file}")
            return False
        
        validation_results['file_exists'] = True
        print(f"✅ Файл презентации существует: {os.path.basename(final_file)}")
        
        # 2. Проверка размера файла
        try:
            file_size = os.path.getsize(final_file)
            # Минимум 100KB для презентации с изображениями (тестовые изображения могут быть маленькими)
            min_size = 100 * 1024
            validation_results['file_size_ok'] = file_size > min_size
            
            if validation_results['file_size_ok']:
                print(f"✅ Размер файла корректный: {file_size / (1024*1024):.2f} MB")
            else:
                print(f"❌ Размер файла слишком мал: {file_size} байт")
                return False
                
        except Exception as e:
            print(f"❌ Ошибка при проверке размера файла: {e}")
            return False
        
        # 3. Проверка структуры презентации
        try:
            prs = Presentation(final_file)
            
            # Проверка количества слайдов
            expected_slides = len(self.slides_data)
            actual_slides = len(prs.slides)
            validation_results['slide_count_correct'] = actual_slides == expected_slides
            
            if validation_results['slide_count_correct']:
                print(f"✅ Количество слайдов корректное: {actual_slides}")
            else:
                print(f"❌ Неправильное количество слайдов: ожидалось {expected_slides}, получено {actual_slides}")
                return False
            
            # Проверка вставки изображений (если AI включен)
            if self.use_ai_illustrations:
                images_found = 0
                ai_images_found = 0
                
                for i, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            images_found += 1
                            # Проверяем, является ли это AI-изображением
                            slide_number = i + 1
                            if self._should_generate_illustration(slide_number, self.slides_data[i]):
                                ai_images_found += 1
                
                expected_ai_images = self.execution_stats.get('images_generated')
                expected_total_images = expected_ai_images + (len(self.template_images) * len(prs.slides))
                
                validation_results['images_inserted'] = ai_images_found >= expected_ai_images
                
                if validation_results['images_inserted']:
                    print(f"✅ AI-изображения вставлены: {ai_images_found}/{expected_ai_images}")
                    print(f"✅ Всего изображений на слайдах: {images_found}")
                else:
                    print(f"❌ Недостаточно AI-изображений: найдено {ai_images_found}, ожидалось {expected_ai_images}")
                    return False
            else:
                validation_results['images_inserted'] = True
                print("✅ AI-изображения не использовались (это ОК)")
            
            validation_results['no_corruption'] = True
            print("✅ Файл презентации читается без ошибок")
            
        except Exception as e:
            print(f"❌ Ошибка при валидации презентации: {e}")
            return False
        
        # Итоговая проверка
        all_passed = all(validation_results.values())
        
        if all_passed:
            print("\n🎉 ВСЕ ПРОВЕРКИ ПРОЙДЕНЫ УСПЕШНО!")
            print("✅ Презентация создана корректно")
            print("✅ Все компоненты на месте")
            print("✅ Качество соответствует стандартам")
        else:
            print("\n❌ ВАЛИДАЦИЯ НЕ ПРОЙДЕНА")
            for check, passed in validation_results.items():
                status = "✅" if passed else "❌"
                print(f"   {status} {check}")
        
        return all_passed

    def _show_final_report(self):
        """Показывает финальный отчет о генерации AI-иллюстраций"""
        stats = self.generation_stats
        
        print(f"\n=== Отчет о генерации AI-иллюстраций ===")
        print(f"Всего слайдов для обработки: {stats['total_slides']}")
        print(f"Промптов сгенерировано: {stats['prompts_generated']}")
        print(f"Промптов с ошибками: {stats['prompts_failed']}")
        print(f"Изображений создано: {stats['images_generated']}")
        print(f"Изображений с ошибками: {stats['images_failed']}")
        
        if stats['total_slides'] > 0:
            success_rate_prompts = (stats['prompts_generated'] / stats['total_slides']) * 100
            success_rate_images = (stats['images_generated'] / stats['total_slides']) * 100
            print(f"Успешность промптов: {success_rate_prompts:.1f}%")
            print(f"Успешность изображений: {success_rate_images:.1f}%")
            
        if self.logger:
            self.logger.info(f"AI generation completed: {json.dumps(stats, ensure_ascii=False)}")
    
    def _remove_shape_borders(self, shape):
        """Удаляет границы и заливку у фигуры"""
        try:
            shape.line.fill.background()  # Прозрачная граница
            shape.fill.background()  # Прозрачная заливка
        except:
            pass  # Игнорируем ошибки с форматированием

    def run(self):
        """Запуск всего процесса генерации"""
        # Красивый ASCII заголовок RW Tech
        ASCIIArt.print_header()
        
        # Анимированный баннер RW Tech
        ColorfulUI.print_rw_tech_banner()
        
        # Анимация загрузки в стиле RW Tech
        ColorfulUI.animated_rw_tech_loading("Инициализация RW Tech системы", 2)
        
        try:
            # ЭТАП 1: Базовая валидация
            self.validate_files()
            self.setup_ai_illustrations()
            
            # ЭТАП 2: Строгая валидация API ключей
            if self.use_ai_illustrations:
                print(f"\n🔒 СТРОГАЯ ВАЛИДАЦИЯ API")
                api_valid = self._validate_and_update_api_keys(interactive=True)
                if not api_valid:
                    print("\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Невалидные API ключи")
                    print("Программа не может продолжить с AI-иллюстрациями")
                    self._cleanup_partial_results()
                    sys.exit(10)  # Код ошибки для невалидных API ключей
            
            # ЭТАП 3: Обработка контента
            self.parse_content()
            self.load_template()
            
            # ЭТАП 4: Генерация базовой презентации
            ColorfulUI.print_rw_tech_step("Создание базовой презентации", "Генерация 60 слайдов с корпоративным дизайном")
            self.generate_presentation()
            
            # ЭТАП 5: Параллельная генерация AI-иллюстраций (критический этап)
            if self.use_ai_illustrations:
                ColorfulUI.print_ascii_step(5, "AI-генерация иллюстраций", "Создание изображений с помощью ИИ")
                ai_success = self._process_ai_illustrations_parallel()
                # _process_ai_illustrations_parallel уже содержит sys.exit() при критических ошибках
            
            # ЭТАП 6: Финальная валидация
            ColorfulUI.print_ascii_step("ВАЛИДАЦИЯ", "Проверка результата", "Комплексная проверка созданной презентации")
            validation_success = self.validate_final_result()
            
            if not validation_success:
                print("\n⛔ КРИТИЧЕСКАЯ ОШИБКА: Финальная валидация не пройдена")
                self._cleanup_partial_results()
                sys.exit(7)  # Код ошибки валидации
            
            # Контрольная точка: финальная валидация
            if not self.checkpoints.validate_checkpoint('final_validation'):
                print("\n⛔ ОСТАНОВКА: Финальная валидация не пройдена")
                sys.exit(8)
            
            # ЭТАП 7: Сохранение истории (при успехе)
            if self.use_ai_illustrations:
                print(f"\n🔸 ЭТАП ФИНАЛ: Сохранение истории")
                self._save_generation_history()
                
                # Очищаем временные файлы ПОСЛЕ сохранения истории
                self._cleanup_partial_results(save_history=False)  # История уже сохранена
            
            # УСПЕХ: Все проверки пройдены
            execution_success = self.execution_stats.print_final_report()
            
            # Красивый баннер успешного завершения RW Tech
            ASCIIArt.print_success_banner()
            print(f"✅ Создано слайдов: {len(self.slides_data)}")
            
            # Определяем итоговый файл
            illustrated_path = os.path.join(self.result_dir, "RWTech_Universal_Presentation_Illustrated.pptx")
            final_file = illustrated_path if (self.use_ai_illustrations and os.path.exists(illustrated_path)) else self.result_file
            
            print(f"✅ Файл сохранен: {final_file}")
            print(f"✅ Размер файла: {os.path.getsize(final_file) / (1024*1024):.2f} MB")
            
            if self.use_ai_illustrations:
                stats = self.execution_stats
                print(f"✅ AI-промпты: {stats.get('prompts_generated')}/{stats.get('prompts_attempted')}")
                print(f"✅ AI-изображения: {stats.get('images_generated')}/{stats.get('images_attempted')}")
                print(f"✅ Вставлено в презентацию: {stats.get('images_inserted')}")
                
                # Показываем историю
                history_base = os.path.join(self.base_path, "history")
                if os.path.exists(history_base):
                    history_dirs = [d for d in os.listdir(history_base) if d.startswith('generation_')]
                    if history_dirs:
                        latest_history = sorted(history_dirs)[-1]
                        print(f"✅ История сохранена в: {os.path.join(history_base, latest_history)}")
                
            if self.logger:
                try:
                    log_files = [f for f in os.listdir(self.logs_dir) if f.endswith('.log')]
                    if log_files:
                        latest_log = sorted(log_files)[-1]
                        print(f"✅ Лог сохранен в: {os.path.join(self.logs_dir, latest_log)}")
                except:
                    pass
            
            if execution_success:
                sys.exit(0)  # Полный успех
            else:
                print("\n⚠️  Программа завершена с предупреждениями")
                sys.exit(9)  # Успех с предупреждениями
            
        except KeyboardInterrupt:
            print()
            print("\033[1;33m")  # Ярко-желтый
            print("""
    ╔════════════════════════════════════════════════════════════════╗
    ║                                                                ║
    ║    ⚠️  ОПЕРАЦИЯ ПРЕРВАНА ПОЛЬЗОВАТЕЛЕМ ⚠️                      ║
    ║                                                                ║
    ║    Работа программы остановлена по запросу пользователя        ║
    ║    Частичные результаты могут быть сохранены                   ║
    ║                                                                ║
    ╚════════════════════════════════════════════════════════════════╝
            """)
            print("\033[0m")
            sys.exit(1)
        except Exception as e:
            print()
            print("\033[1;31m")  # Ярко-красный
            print("""
    ╔════════════════════════════════════════════════════════════════╗
    ║                                                                ║
    ║    ❌ КРИТИЧЕСКАЯ ОШИБКА ВЫПОЛНЕНИЯ ❌                        ║
    ║                                                                ║
    ║    Произошла неожиданная ошибка в работе программы             ║
    ║    Проверьте логи для получения подробной информации           ║
    ║                                                                ║
    ╚════════════════════════════════════════════════════════════════╝
            """)
            print("\033[0m")
            print(f"Детали ошибки: {e}")
            sys.exit(1)


if __name__ == "__main__":
    generator = RWTechPPTXGenerator()
    generator.run()