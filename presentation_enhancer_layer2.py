"""
Presentation Enhancer Layer 2 - Title optimization and layout decoration
"""

import json
import time
import re
from typing import Dict, List, Tuple
from anthropic import Anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class PresentationEnhancerLayer2:
    def __init__(self, api_key: str):
        self.client = Anthropic(api_key=api_key)
        self.model = "claude-3-5-sonnet-20241022"
        
    def enhance_layer2(self, presentation_path: str, output_path: str) -> Dict:
        """
        Применяет второй слой улучшений: оптимизация заголовков и декорация макета
        """
        print("🎯 Запуск второго слоя улучшений...")
        
        # Загружаем презентацию
        prs = Presentation(presentation_path)
        results = {
            "total_slides": len(prs.slides),
            "optimized_titles": 0,
            "optimized_text": 0,
            "decorated_layouts": 0,
            "errors": []
        }
        
        # Шаг 1: Оптимизация заголовков и текста через Claude
        print("📝 Оптимизация заголовков и текста...")
        self._optimize_content_with_claude(prs, results)
        
        # Шаг 2: Декорация макета (уменьшение ширины и выравнивание)
        print("🎨 Применение декорации макета...")
        self._apply_layout_decoration(prs, results)
        
        # Сохраняем результат
        prs.save(output_path)
        
        print(f"✅ Второй слой применен! Результат сохранен: {output_path}")
        return results
    
    def _optimize_content_with_claude(self, prs: Presentation, results: Dict):
        """Оптимизирует заголовки и текст с помощью Claude"""
        
        for i, slide in enumerate(prs.slides):
            try:
                print(f"🔄 Обработка слайда {i+1}/{len(prs.slides)}...")
                
                # Находим заголовок и основной текст
                title_shape = None
                content_shape = None
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if shape.placeholder_format and shape.placeholder_format.type == 1:  # Title
                            title_shape = shape
                        elif shape.placeholder_format and shape.placeholder_format.type == 2:  # Content
                            content_shape = shape
                        elif not title_shape and len(shape.text.strip()) < 100:  # Короткий текст - вероятно заголовок
                            title_shape = shape
                        elif not content_shape:  # Длинный текст - вероятно контент
                            content_shape = shape
                
                if title_shape or content_shape:
                    # Получаем оптимизированный контент от Claude
                    optimized = self._get_optimized_content(
                        title_shape.text if title_shape else "",
                        content_shape.text if content_shape else ""
                    )
                    
                    # Применяем оптимизированный заголовок с жирным шрифтом
                    if title_shape and optimized.get('title'):
                        title_shape.text = optimized['title']
                        # Делаем заголовок жирным
                        if hasattr(title_shape, 'text_frame'):
                            for paragraph in title_shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True
                        results["optimized_titles"] += 1
                    
                    # Применяем оптимизированный текст
                    if content_shape and optimized.get('content'):
                        content_shape.text = optimized['content']
                        # Принудительно применяем форматирование после установки текста
                        if hasattr(content_shape, 'text_frame'):
                            for paragraph in content_shape.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.LEFT  # Выравнивание по левому краю
                                for run in paragraph.runs:
                                    run.font.size = Pt(20)  # 20pt размер шрифта
                        results["optimized_text"] += 1
                        
            except Exception as e:
                results["errors"].append(f"Слайд {i+1}: {str(e)}")
                print(f"⚠️ Ошибка при обработке слайда {i+1}: {str(e)}")
    
    def _get_optimized_content(self, title: str, content: str) -> Dict[str, str]:
        """Получает оптимизированный контент от Claude"""
        
        prompt = f"""
Ты эксперт по созданию качественных презентаций для руководителей Гознака.

ЗАДАЧА: Оптимизировать заголовок и текст слайда согласно требованиям.

ИСХОДНЫЕ ДАННЫЕ:
Заголовок: {title}
Текст: {content}

ТРЕБОВАНИЯ К ЗАГОЛОВКУ:
1. Длина: строго 5-6 слов
2. Формат: первые 3 слова на первой строке, оставшиеся 2-3 слова на второй строке
3. Должен отражать основную мысль слайда
4. Качественный и профессиональный

ТРЕБОВАНИЯ К ТЕКСТУ:
1. Максимум 1 абзац
2. Лаконичный и интересный
3. Профессиональный стиль для руководителей
4. Сохранить ключевую информацию

ФОРМАТ ОТВЕТА (только JSON):
{{
    "title": "Первые три слова\\nОставшиеся слова",
    "content": "Оптимизированный текст в один абзац"
}}
"""
        
        try:
            response = self.client.messages.create(
                model=self.model,
                max_tokens=4000,
                temperature=0.7,
                messages=[{"role": "user", "content": prompt}]
            )
            
            # Извлекаем JSON из ответа
            response_text = response.content[0].text
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            
            if json_match:
                return json.loads(json_match.group())
            else:
                return {"title": title, "content": content}
                
        except Exception as e:
            print(f"⚠️ Ошибка Claude API: {str(e)}")
            return {"title": title, "content": content}
    
    def _apply_layout_decoration(self, prs: Presentation, results: Dict):
        """Применяет декорацию макета: уменьшение ширины и выравнивание только основного текста"""
        
        for i, slide in enumerate(prs.slides):
            try:
                # Находим заголовок, основной текст и изображения отдельно
                title_shape = None
                content_shapes = []
                image_shapes = []
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text.strip():
                        # Определяем заголовок
                        if shape.placeholder_format and shape.placeholder_format.type == 1:  # Title placeholder
                            title_shape = shape
                        # Определяем основной контент
                        elif shape.placeholder_format and shape.placeholder_format.type == 2:  # Content placeholder
                            content_shapes.append(shape)
                        # Если нет placeholder_format, определяем по размеру текста и позиции
                        elif not title_shape and len(shape.text.strip()) < 100 and shape.top < prs.slide_height // 3:
                            title_shape = shape  # Короткий текст в верхней части - заголовок
                        elif shape != title_shape:  # Остальные блоки - контент
                            content_shapes.append(shape)
                    elif hasattr(shape, 'image'):  # Изображения
                        image_shapes.append(shape)
                
                # Применяем декорацию только к основному тексту (НЕ к заголовку)
                if content_shapes:
                    for shape in content_shapes:
                        # Уменьшаем ширину на 40%
                        original_width = shape.width
                        new_width = int(original_width * 0.6)  # Уменьшаем на 40%
                        shape.width = new_width
                        
                        # Определяем выравнивание контейнера (через один справа-слева)
                        if i % 2 == 1:  # Слайды 2, 4, 6... - контейнер справа
                            # Сдвигаем контейнер вправо
                            slide_width = prs.slide_width
                            shape.left = slide_width - shape.width - Inches(0.5)
                        else:  # Слайды 1, 3, 5... - контейнер слева
                            # Сдвигаем контейнер влево
                            shape.left = Inches(0.5)
                        
                        # Настройка текста: выравнивание по левому краю и размер шрифта
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.LEFT  # Выравнивание по левому краю
                                for run in paragraph.runs:
                                    run.font.size = Pt(20)  # 20pt размер шрифта
                        
                        # Позиционирование: на 15% выше центра доступного пространства
                        title_bottom = title_shape.top + title_shape.height if title_shape else Inches(1.5)
                        available_height = prs.slide_height - title_bottom - Inches(0.5)  # Отступ снизу
                        center_position = title_bottom + (available_height - shape.height) // 2
                        # Поднимаем на 15% от центральной позиции
                        offset_up = available_height * 0.15
                        shape.top = int(center_position - offset_up)
                
                # Обработка изображений: зеркальное отражение на слайдах с правым выравниванием
                if image_shapes:
                    self._handle_images(slide, image_shapes, i, prs)
                    
                results["decorated_layouts"] += 1
                    
            except Exception as e:
                results["errors"].append(f"Декорация слайда {i+1}: {str(e)}")
                print(f"⚠️ Ошибка декорации слайда {i+1}: {str(e)}")
    
    def _handle_images(self, slide, image_shapes, slide_index, prs):
        """Обрабатывает изображения: зеркальное отражение на слайдах с правым выравниванием"""
        
        try:
            # Находим изображения в правом нижнем углу
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            right_bottom_images = []
            for shape in image_shapes:
                # Проверяем, что изображение в правой части слайда (больше 50% ширины)
                # и в нижней части (больше 50% высоты)
                if (shape.left + shape.width/2 > slide_width/2 and 
                    shape.top + shape.height/2 > slide_height/2):
                    right_bottom_images.append(shape)
            
            if not right_bottom_images:
                return
            
            # Определяем, куда двигать изображения
            if slide_index % 2 == 1:  # Слайды 2, 4, 6... - текст справа, изображение переносим влево
                for img_shape in right_bottom_images:
                    # Зеркально отражаем позицию: из правого угла в левый
                    new_left = slide_width - (img_shape.left + img_shape.width)
                    img_shape.left = new_left
                    
                    # Удаляем оригинал из правой части (создаем копию в левой части)
                    # Оригинал остается, но мы его переместили влево
            
            # На слайдах 1, 3, 5... (текст слева) - изображения остаются справа (ничего не делаем)
                    
        except Exception as e:
            print(f"⚠️ Ошибка обработки изображений на слайде {slide_index+1}: {str(e)}")


def main():
    """Тестирование второго слоя улучшений"""
    import os
    
    # API ключ
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        print("❌ Не найден ANTHROPIC_API_KEY")
        return
    
    # Создаем улучшатор
    enhancer = PresentationEnhancerLayer2(api_key)
    
    # Тестируем на последней сгенерированной презентации
    input_file = "result/presentation.pptx"
    output_file = "result/enhanced_layer2_presentation.pptx"
    
    if os.path.exists(input_file):
        results = enhancer.enhance_layer2(input_file, output_file)
        
        print("\n📊 РЕЗУЛЬТАТЫ ВТОРОГО СЛОЯ:")
        print(f"Всего слайдов: {results['total_slides']}")
        print(f"Оптимизировано заголовков: {results['optimized_titles']}")
        print(f"Оптимизировано текстов: {results['optimized_text']}")
        print(f"Декорировано макетов: {results['decorated_layouts']}")
        
        if results['errors']:
            print(f"\n⚠️ Ошибки ({len(results['errors'])}):")
            for error in results['errors']:
                print(f"  - {error}")
    else:
        print(f"❌ Файл {input_file} не найден")


if __name__ == "__main__":
    main()