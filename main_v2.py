#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ADHOC PowerPoint Generator v2 - с красивым дизайном
Использует готовый шаблон с сохранением всех элементов дизайна
"""

import os
import sys
import re
from typing import List, Tuple, Dict
from pathlib import Path
import copy

try:
    from striprtf.striprtf import rtf_to_text
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"ОШИБКА: Не установлены зависимости: {e}")
    print("Запустите: pip install python-pptx striprtf")
    sys.exit(1)


class TemplateBasedConverter:
    """Конвертер RTF в PowerPoint на основе готового шаблона"""
    
    def __init__(self, rtf_path: str, template_path: str, output_path: str):
        self.rtf_path = Path(rtf_path)
        self.template_path = Path(template_path)
        self.output_path = Path(output_path)
        self.max_content_length = 350  # Максимум символов на слайд
        self.min_content_for_slide = 60  # Минимум контента для отдельного слайда
        
        # Проверяем входные файлы
        if not self.rtf_path.exists():
            raise FileNotFoundError(f"RTF файл не найден: {self.rtf_path}")
        if not self.template_path.exists():
            raise FileNotFoundError(f"Шаблон не найден: {self.template_path}")
    
    def parse_rtf(self) -> List[Tuple[str, List[str]]]:
        """Парсит RTF файл и возвращает структурированные данные"""
        print(f"📄 Читаем RTF файл: {self.rtf_path}")
        
        try:
            # Читаем RTF с правильной кодировкой
            with open(self.rtf_path, 'r', encoding='cp1251') as f:
                rtf_content = f.read()
            
            # Конвертируем в текст
            text = rtf_to_text(rtf_content)
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            
            print(f"✅ Прочитано {len(lines)} строк, {len(text)} символов")
            
            # Умная группировка контента
            structured_data = []
            current_title = None
            current_content = []
            content_length = 0
            
            for line in lines:
                if self._is_heading(line):
                    # Сохраняем предыдущий блок только если есть достаточно контента
                    if current_title and content_length >= self.min_content_for_slide:
                        structured_data.append((current_title, current_content))
                        current_title = line
                        current_content = []
                        content_length = 0
                    elif current_title:
                        # Объединяем с предыдущим заголовком если мало контента
                        current_title = f"{current_title} | {line}"
                    else:
                        # Первый заголовок
                        current_title = line
                        current_content = []
                        content_length = 0
                else:
                    # Добавляем к текущему контенту
                    current_content.append(line)
                    content_length += len(line)
            
            # Не забываем последний блок
            if current_title and current_content:
                structured_data.append((current_title, current_content))
            
            print(f"📊 Найдено {len(structured_data)} секций")
            return structured_data
            
        except UnicodeDecodeError:
            print("⚠️  Ошибка кодировки, пробуем UTF-8")
            try:
                with open(self.rtf_path, 'r', encoding='utf-8') as f:
                    rtf_content = f.read()
                text = rtf_to_text(rtf_content)
                return self._process_text_to_structure(text)
            except Exception as e:
                raise Exception(f"Не удалось прочитать RTF файл: {e}")
        
        except Exception as e:
            raise Exception(f"Ошибка при парсинге RTF: {e}")
    
    def _is_heading(self, line: str) -> bool:
        """Определяет, является ли строка заголовком"""
        # Исключаем очень короткие строки
        if len(line) < 15:
            return False
        
        # Исключаем очень длинные строки  
        if len(line) > 100:
            return False
            
        # Исключаем строки состоящие только из цифр/спецсимволов/валют
        if re.match(r'^[\d\s\$%№€₽-]+$', line):
            return False
            
        # Исключаем строки начинающиеся с цифр (статистика)
        if re.match(r'^\d+[%\d\s\$€₽]', line):
            return False
            
        # Признаки заголовка:
        # 1. Заканчивается двоеточием
        if line.endswith(':'):
            return True
            
        # 2. Содержит ключевые слова заголовков
        heading_words = ['эра', 'революция', 'эволюция', 'введение', 'что такое', 
                        'как', 'почему', 'зачем', 'вопрос', 'проблема', 'решение']
        if any(word in line.lower() for word in heading_words):
            return True
            
        # 3. Средняя длина и нет точки в конце (не предложение)
        if 20 <= len(line) <= 80 and not line.endswith('.'):
            return True
            
        return False
    
    def _process_content_for_slides(self, title: str, content: List[str]) -> List[Tuple[str, str]]:
        """Обрабатывает контент для создания слайдов"""
        if not content:
            return [(title, "")]
        
        # Объединяем весь контент
        full_content = "\n".join(content)
        
        # Если контент помещается на один слайд
        if len(full_content) <= self.max_content_length:
            return [(title, self._format_content(full_content))]
        
        # Разбиваем на части
        parts = self._split_content_smartly(content)
        slides = []
        
        for i, part in enumerate(parts, 1):
            if len(parts) > 1:
                part_title = f"{title} (часть {i})"
            else:
                part_title = title
            slides.append((part_title, self._format_content(part)))
        
        return slides
    
    def _split_content_smartly(self, content: List[str]) -> List[str]:
        """Умно разбивает контент на части, сохраняя смысл"""
        parts = []
        current_part = []
        current_length = 0
        
        for line in content:
            line_length = len(line)
            
            # Если добавление этой строки превысит лимит
            if current_length + line_length > self.max_content_length and current_part:
                # Сохраняем текущую часть
                parts.append("\n".join(current_part))
                current_part = [line]
                current_length = line_length
            else:
                current_part.append(line)
                current_length += line_length + 1  # +1 для переноса строки
        
        # Добавляем последнюю часть
        if current_part:
            parts.append("\n".join(current_part))
        
        return parts
    
    def _format_content(self, content: str) -> str:
        """Форматирует контент для слайда"""
        if not content:
            return ""
        
        lines = content.split('\n')
        formatted_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Если строка не очень длинная, делаем её bullet point
            if len(line) <= 200:
                if not line.startswith('•') and not line.startswith('-'):
                    line = f"• {line}"
            
            formatted_lines.append(line)
        
        return '\n'.join(formatted_lines)
    
    def _create_slide_with_template_design(self, presentation, template_presentation):
        """Создает слайд используя шаблон как основу для новой презентации"""
        # Просто используем макет Title and Content, но применим дизайн шаблона
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content
        
        # Применяем цвета и стили из шаблона
        self._apply_template_styling(slide, template_presentation)
        
        return slide
    
    def _apply_template_styling(self, slide, template_presentation):
        """Применяет стили из шаблона к слайду"""
        try:
            # Получаем цветовую схему из шаблона
            template_slide = template_presentation.slides[0]
            
            # Устанавливаем фон слайда если есть
            if hasattr(template_slide, 'background'):
                slide.background = template_slide.background
                
        except Exception as e:
            print(f"⚠️  Не удалось применить все стили шаблона: {e}")
            # Продолжаем работу без стилей
    
    def _update_slide_content(self, slide, title: str, content: str):
        """Обновляет содержимое слайда используя стандартные placeholders"""
        try:
            # Используем стандартные placeholders для Title and Content макета
            if slide.shapes.title:
                slide.shapes.title.text = title
                self._format_title(slide.shapes.title)
            
            # Ищем content placeholder
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    shape.text = content
                    self._format_content_shape(shape)
                    break
            else:
                # Если не нашли placeholder, ищем по тексту
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape != slide.shapes.title:
                        shape.text = content
                        self._format_content_shape(shape)
                        break
                        
        except Exception as e:
            print(f"⚠️  Проблема с обновлением контента слайда: {e}")
            # Fallback - просто используем первые два текстовых элемента
            text_shapes = [s for s in slide.shapes if hasattr(s, 'text')]
            if len(text_shapes) >= 1:
                text_shapes[0].text = title
                self._format_title(text_shapes[0])
            if len(text_shapes) >= 2:
                text_shapes[1].text = content
                self._format_content_shape(text_shapes[1])
    
    def _format_title(self, shape):
        """Форматирует заголовок: шрифт Montserrat 30pt"""
        if not hasattr(shape, 'text_frame'):
            return
            
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = 'Montserrat'
                run.font.size = Pt(30)
                run.font.bold = True
    
    def _format_content_shape(self, shape):
        """Форматирует контент: шрифт Montserrat 18pt"""
        if not hasattr(shape, 'text_frame'):
            return
            
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = 'Montserrat'
                run.font.size = Pt(18)
                run.font.bold = False
    
    def create_presentation(self, structured_data: List[Tuple[str, List[str]]]) -> None:
        """Создает PowerPoint презентацию на основе шаблона"""
        print("🎨 Создаем PowerPoint презентацию на основе шаблона...")
        
        try:
            # Загружаем шаблон как основу для новой презентации
            template_prs = Presentation(str(self.template_path))
            
            # Создаем новую презентацию НА ОСНОВЕ шаблона
            prs = Presentation(str(self.template_path))
            
            # Удаляем все слайды из шаблона, оставляем только структуру
            slide_ids_to_remove = list(range(len(prs.slides)))
            for slide_id in reversed(slide_ids_to_remove):
                rId = prs.slides._sldIdLst[slide_id].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[slide_id]
            
            slide_count = 0
            
            # Добавляем титульный слайд
            title_slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
            self._update_slide_content(title_slide, 
                                     "Презентация на основе RTF документа", 
                                     f"Автоматически создано из {self.rtf_path.name}")
            slide_count += 1
            
            # Обрабатываем каждую секцию
            for title, content in structured_data:
                slides_for_section = self._process_content_for_slides(title, content)
                
                for slide_title, slide_content in slides_for_section:
                    # Создаем слайд с дизайном шаблона
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
                    self._update_slide_content(slide, slide_title, slide_content)
                    slide_count += 1
            
            print(f"✅ Создано {slide_count} слайдов")
            
            # Сохраняем презентацию
            print(f"💾 Сохраняем в: {self.output_path}")
            prs.save(str(self.output_path))
            print(f"🎉 Презентация готова: {self.output_path}")
            
        except Exception as e:
            raise Exception(f"Ошибка при создании презентации: {e}")
    
    def convert(self) -> None:
        """Основной метод конвертации"""
        print("🚀 Запуск конвертации RTF → PowerPoint (v2 с дизайном)")
        print("=" * 60)
        
        try:
            # Парсим RTF
            structured_data = self.parse_rtf()
            
            # Создаем презентацию
            self.create_presentation(structured_data)
            
            print("=" * 60)
            print("✅ КОНВЕРТАЦИЯ ЗАВЕРШЕНА УСПЕШНО!")
            
        except Exception as e:
            print(f"❌ ОШИБКА: {e}")
            import traceback
            traceback.print_exc()
            sys.exit(1)


def main():
    """Главная функция"""
    # Пути к файлам
    script_dir = Path(__file__).parent
    rtf_path = script_dir / "content" / "content1"
    template_path = script_dir / "tempate" / "Шаблон презентации 16х9.pptx"
    output_path = script_dir / "presentation_v2.pptx"
    
    print("🎯 ADHOC PowerPoint Generator v2 (С ДИЗАЙНОМ)")
    print(f"📂 RTF файл: {rtf_path}")
    print(f"📂 Шаблон: {template_path}")
    print(f"📂 Выходной файл: {output_path}")
    print()
    
    # Создаем конвертер и запускаем
    converter = TemplateBasedConverter(str(rtf_path), str(template_path), str(output_path))
    converter.convert()


if __name__ == "__main__":
    main()